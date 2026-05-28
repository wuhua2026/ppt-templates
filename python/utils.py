"""PPT模板工具函数"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math


def hex_to_rgb(hex_color):
    """十六进制颜色转RGB，也支持直接传入RGBColor对象"""
    if isinstance(hex_color, RGBColor):
        return hex_color
    hex_color = str(hex_color).lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def create_circle_positions(center_x, center_y, radius, count, start_angle=0):
    """计算圆周上均匀分布的位置"""
    positions = []
    for i in range(count):
        angle = start_angle + (2 * math.pi * i / count)
        x = center_x + radius * math.cos(angle)
        y = center_y + radius * math.sin(angle)
        positions.append((x, y))
    return positions


def create_grid_positions(rows, cols, start_x, start_y, cell_width, cell_height, gap=0):
    """计算网格布局位置"""
    positions = []
    for row in range(rows):
        for col in range(cols):
            x = start_x + col * (cell_width + gap)
            y = start_y + row * (cell_height + gap)
            positions.append((x, y))
    return positions


def add_shadow(shape, blur=Pt(5), offset=Pt(3), color=RGBColor(0, 0, 0), opacity=0.3):
    """为形状添加阴影效果"""
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    spPr = shape._element.find('.//a:spPr', nsmap)
    if spPr is None:
        spPr = etree.SubElement(shape._element, '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
    effectLst = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
    outerShdw = etree.SubElement(effectLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw')
    outerShdw.set('blurRad', str(int(blur)))
    outerShdw.set('dist', str(int(offset)))
    outerShdw.set('dir', '5400000')
    outerShdw.set('algn', 'tl')
    srgbClr = etree.SubElement(outerShdw, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}' if isinstance(color, tuple) else str(color))
    alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
    alpha.set('val', str(int(opacity * 100000)))


def set_shape_transparency(shape, transparency):
    """设置形状透明度 (0-1)"""
    from lxml import etree
    solidFill = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', str(int((1 - transparency) * 100000)))


def format_text(text_frame, text, font_size=18, font_color=None, bold=False,
                alignment=PP_ALIGN.LEFT, font_name=None):
    """格式化文本框内容"""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color
    if font_name:
        run.font.name = font_name
    p.alignment = alignment
    text_frame.word_wrap = True
    return text_frame

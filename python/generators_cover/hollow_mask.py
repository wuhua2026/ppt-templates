"""镂空遮罩封面生成器 - 几何形状镂空效果，层次感强"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency


class HollowMaskCover(TemplateGenerator):
    """镂空遮罩封面"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._title_font_size = 42
        self._subtitle_font_size = 18
        self._title_color = RGBColor(255, 255, 255)
        self._subtitle_color = RGBColor(200, 200, 200)
        self._mask_color = None
        self._overlay_opacity = 0.55

    def set_title(self, title, font_size=None, color=None):
        """设置标题"""
        self._title = title
        if font_size is not None:
            self._title_font_size = font_size
        if color is not None:
            self._title_color = color
        return self

    def set_subtitle(self, subtitle, font_size=None, color=None):
        """设置副标题"""
        self._subtitle = subtitle
        if font_size is not None:
            self._subtitle_font_size = font_size
        if color is not None:
            self._subtitle_color = color
        return self

    def _get_theme_colors(self):
        """获取主题颜色配置"""
        if self.theme:
            primary = hex_to_rgb(self.theme.primary)
            secondary = hex_to_rgb(self.theme.secondary)
            accent = hex_to_rgb(self.theme.accent)
        else:
            primary = RGBColor(15, 15, 35)
            secondary = RGBColor(26, 26, 62)
            accent = RGBColor(0, 210, 255)
        return primary, secondary, accent

    def generate(self):
        """生成镂空遮罩封面"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height
        primary, secondary, accent = self._get_theme_colors()

        if self._mask_color is None:
            self._mask_color = primary

        # --- 底层渐变背景 ---
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), slide_w, slide_h,
            [primary, secondary], angle=135
        )

        # --- 大号镂空圆形（右上方） ---
        circle_size = Inches(6.0)
        self.add_shape(
            MSO_SHAPE.OVAL,
            slide_w - circle_size + Inches(0.3),
            Inches(-2.0),
            circle_size, circle_size,
            fill_color=accent
        )
        set_shape_transparency(self.slide.shapes[-1], 0.88)

        # --- 中号镂空圆形（左下方） ---
        circle_mid = Inches(4.0)
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1.5),
            slide_h - circle_mid + Inches(0.5),
            circle_mid, circle_mid,
            fill_color=accent
        )
        set_shape_transparency(self.slide.shapes[-1], 0.92)

        # --- 半透明矩形遮罩层（文字背景） ---
        mask_w = Inches(5.5)
        mask_h = Inches(2.8)
        mask_x = (slide_w - mask_w) // 2
        mask_y = (slide_h - mask_h) // 2

        mask_shape = self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            mask_x, mask_y, mask_w, mask_h,
            fill_color=RGBColor(0, 0, 0)
        )
        set_shape_transparency(mask_shape, 0.45)
        mask_shape.line.fill.background()

        # --- 小菱形装饰 ---
        diamond_size = Inches(0.7)
        self.add_shape(
            MSO_SHAPE.DIAMOND,
            mask_x - Inches(0.9),
            mask_y + mask_h // 2 - diamond_size // 2,
            diamond_size, diamond_size,
            fill_color=accent
        )
        set_shape_transparency(self.slide.shapes[-1], 0.3)

        # --- 标题 ---
        if self._title:
            self.add_textbox(
                self._title,
                mask_x + Inches(0.5), mask_y + Inches(0.4),
                mask_w - Inches(1.0), Inches(1.2),
                font_size=self._title_font_size,
                font_color=self._title_color,
                bold=True,
                alignment=PP_ALIGN.CENTER
            )

        # --- 分隔细线 ---
        line_w = Inches(2.0)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (slide_w - line_w) // 2,
            mask_y + mask_h // 2 + Inches(0.1),
            line_w, Pt(2),
            fill_color=accent
        )

        # --- 副标题 ---
        if self._subtitle:
            self.add_textbox(
                self._subtitle,
                mask_x + Inches(0.5), mask_y + mask_h // 2 + Inches(0.4),
                mask_w - Inches(1.0), Inches(0.8),
                font_size=self._subtitle_font_size,
                font_color=self._subtitle_color,
                bold=False,
                alignment=PP_ALIGN.CENTER
            )

        return self


if __name__ == "__main__":
    cover = HollowMaskCover()
    cover.set_title("镂空遮罩封面", font_size=44)
    cover.set_subtitle("Hollow Mask Cover", font_size=20)
    cover.save("output/hollow_mask_cover.pptx")
    print("镂空遮罩封面已生成: output/hollow_mask_cover.pptx")

"""分屏封面生成器 - 左右分屏设计，色彩与留白的对比"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb


class SplitScreenCover(TemplateGenerator):
    """分屏封面"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._title_font_size = 40
        self._subtitle_font_size = 18
        self._left_color = None
        self._right_color = None
        self._title_color = RGBColor(255, 255, 255)
        self._subtitle_color = RGBColor(60, 60, 60)
        self._split_ratio = 0.5  # 0.0~1.0 左侧占比

    def set_title(self, title, font_size=None, color=None):
        """设置标题（显示在左侧彩色区域）"""
        self._title = title
        if font_size is not None:
            self._title_font_size = font_size
        if color is not None:
            self._title_color = color
        return self

    def set_subtitle(self, subtitle, font_size=None, color=None):
        """设置副标题（显示在右侧白色区域）"""
        self._subtitle = subtitle
        if font_size is not None:
            self._subtitle_font_size = font_size
        if color is not None:
            self._subtitle_color = color
        return self

    def set_split_ratio(self, ratio):
        """设置分屏比例 (0.0~1.0)，左侧占比"""
        self._split_ratio = max(0.1, min(0.9, ratio))
        return self

    def _get_theme_colors(self):
        """获取主题颜色配置"""
        if self.theme:
            primary = hex_to_rgb(self.theme.primary)
            secondary = hex_to_rgb(self.theme.secondary)
            accent = hex_to_rgb(self.theme.accent)
        else:
            primary = RGBColor(44, 62, 80)
            secondary = RGBColor(52, 152, 219)
            accent = RGBColor(231, 76, 60)
        return primary, secondary, accent

    def generate(self):
        """生成分屏封面"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height
        primary, secondary, accent = self._get_theme_colors()

        if self._left_color is None:
            self._left_color = primary
        if self._right_color is None:
            self._right_color = RGBColor(250, 250, 250)

        split_x = int(slide_w * self._split_ratio)
        right_w = slide_w - split_x

        # --- 左侧：渐变色块 ---
        left_bg = self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), split_x, slide_h,
            [primary, secondary], angle=135
        )

        # --- 右侧：浅色背景 ---
        right_bg = self.add_shape(
            MSO_SHAPE.RECTANGLE,
            split_x, Inches(0), right_w, slide_h,
            fill_color=self._right_color
        )

        # --- 垂直分隔线 ---
        line_color = accent
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            split_x - Pt(1), Inches(0),
            Pt(3), slide_h,
            fill_color=line_color
        )

        # --- 左侧几何装饰 ---
        deco_size = Inches(1.8)
        self.add_shape(
            MSO_SHAPE.OVAL,
            split_x - deco_size - Inches(0.3),
            Inches(0.5),
            deco_size, deco_size,
            fill_color=accent
        )
        from python.utils import set_shape_transparency
        set_shape_transparency(self.slide.shapes[-1], 0.75)

        # --- 左侧标题 ---
        margin = Inches(0.8)
        self.add_textbox(
            self._title,
            margin, slide_h // 2 - Inches(0.6),
            split_x - margin * 2, Inches(1.2),
            font_size=self._title_font_size,
            font_color=self._title_color,
            bold=True,
            alignment=PP_ALIGN.LEFT
        )

        # --- 右侧副标题 ---
        self.add_textbox(
            self._subtitle,
            split_x + Inches(0.6),
            slide_h // 2 - Inches(0.4),
            right_w - Inches(1.2), Inches(0.8),
            font_size=self._subtitle_font_size,
            font_color=self._subtitle_color,
            bold=False,
            alignment=PP_ALIGN.LEFT
        )

        # --- 右下角装饰小方块 ---
        block_size = Inches(0.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            split_x + right_w - block_size - Inches(0.6),
            slide_h - block_size - Inches(0.6),
            block_size, block_size,
            fill_color=secondary
        )
        set_shape_transparency(self.slide.shapes[-1], 0.6)

        return self


if __name__ == "__main__":
    cover = SplitScreenCover()
    cover.set_title("分屏封面", font_size=44)
    cover.set_subtitle("Split Screen Cover", font_size=20)
    cover.save("output/split_screen_cover.pptx")
    print("分屏封面已生成: output/split_screen_cover.pptx")

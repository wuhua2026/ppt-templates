"""PPT模板生成器基础类"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree


class TemplateGenerator:
    """所有模板生成器的基类"""

    def __init__(self, theme=None, slide_width=Inches(10), slide_height=Inches(7.5)):
        self.prs = Presentation()
        self.prs.slide_width = slide_width
        self.prs.slide_height = slide_height
        self.theme = theme
        self.slide = None
        self._shape_id_counter = 1000

    def create_slide(self):
        """创建空白幻灯片"""
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        return self.slide

    def add_shape(self, shape_type, left, top, width, height, fill_color=None, outline_color=None, outline_width=None):
        """添加形状"""
        shape = self.slide.shapes.add_shape(shape_type, left, top, width, height)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if outline_color:
            shape.line.color.rgb = outline_color
            if outline_width:
                shape.line.width = outline_width
        else:
            shape.line.fill.background()
        return shape

    def add_gradient_shape(self, shape_type, left, top, width, height, colors, angle=0):
        """添加渐变填充形状"""
        shape = self.slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.gradient()
        shape.fill.gradient_stops[0].color.rgb = colors[0]
        shape.fill.gradient_stops[0].position = 0.0
        if len(colors) > 1:
            shape.fill.gradient_stops[1].color.rgb = colors[1]
            shape.fill.gradient_stops[1].position = 1.0
        shape.line.fill.background()
        return shape

    def add_textbox(self, text, left, top, width, height, font_size=18, font_color=None,
                    bold=False, alignment=PP_ALIGN.LEFT, font_name=None, line_spacing=1.2):
        """添加文本框"""
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        if font_color:
            p.font.color.rgb = font_color
        if font_name:
            p.font.name = font_name
        return txBox

    def add_paragraph(self, text_frame, text, font_size=14, font_color=None,
                      bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
        """向已有文本框添加段落"""
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        if font_color:
            p.font.color.rgb = font_color
        if font_name:
            p.font.name = font_name
        return p

    def add_image(self, image_path, left, top, width=None, height=None):
        """添加图片"""
        if width and height:
            return self.slide.shapes.add_picture(image_path, left, top, width, height)
        elif width:
            return self.slide.shapes.add_picture(image_path, left, top, width=width)
        elif height:
            return self.slide.shapes.add_picture(image_path, left, top, height=height)
        return self.slide.shapes.add_picture(image_path, left, top)

    def add_transition(self, transition_type="fade", speed="med", duration=5000):
        """添加幻灯片转场效果（通过XML操作）"""
        transition_map = {
            "fade": "fade",
            "push": "push",
            "cover": "cover",
            "wipe": "wipe",
            "split": "split",
            "blinds": "blinds",
            "cut": "cut",
            "dissolve": "dissolve",
        }
        speed_map = {"slow": 2000, "med": 1000, "fast": 500}
        adv_tm = speed_map.get(speed, 1000)

        xml_str = f'''<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            advClick="1" advTm="{adv_tm}">
            <p:{transition_type} spd="med" dir="horz"/>
        </p:transition>'''
        self.slide._element.append(etree.fromstring(xml_str))

    def _get_next_shape_id(self):
        """获取下一个形状ID"""
        self._shape_id_counter += 1
        return self._shape_id_counter

    def save(self, filepath):
        """保存PPT文件"""
        self.prs.save(filepath)
        return filepath

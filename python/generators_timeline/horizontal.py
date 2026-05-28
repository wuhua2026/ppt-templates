"""水平时间轴 - 经典水平线时间轴，节点均匀分布，标签上下交替"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb


class HorizontalTimeline(TemplateGenerator):
    """水平时间轴"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._events = []  # list of dicts: {label, date, detail, color}
        self._node_radius = Inches(0.15)
        self._line_color = RGBColor(189, 195, 199)
        self._default_colors = [
            RGBColor(52, 152, 219),
            RGBColor(231, 76, 60),
            RGBColor(46, 204, 113),
            RGBColor(155, 89, 182),
            RGBColor(241, 196, 15),
            RGBColor(230, 126, 34),
        ]

    def set_title(self, title, subtitle=""):
        self._title = title
        self._subtitle = subtitle
        return self

    def set_data(self, events):
        """设置事件列表

        Args:
            events: list of dict, 每项包含:
                label (str): 事件名称
                date (str): 日期
                detail (str): 详细描述
                color (RGBColor, optional): 自定义颜色
        """
        self._events = events
        return self

    def generate(self):
        """生成水平时间轴"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # 标题
        if self._title:
            self.add_textbox(
                self._title, Inches(0.5), Inches(0.3),
                slide_w - Inches(1), Inches(0.6),
                font_size=28, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.LEFT
            )
        if self._subtitle:
            self.add_textbox(
                self._subtitle, Inches(0.5), Inches(0.85),
                slide_w - Inches(1), Inches(0.4),
                font_size=14,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.LEFT
            )

        # 水平线位置
        line_y = int(slide_h * 0.55)
        margin_x = Inches(0.8)
        line_start_x = int(margin_x)
        line_end_x = int(slide_w - margin_x)

        # 绘制水平主轴线
        line = self.slide.shapes.add_connector(
            1,  # straight connector
            line_start_x, line_y,
            line_end_x, line_y
        )
        line.line.color.rgb = self._line_color
        line.line.width = Pt(3)

        n = len(self._events)
        if n == 0:
            return self

        spacing = (line_end_x - line_start_x) / max(n - 1, 1)

        for i, event in enumerate(self._events):
            if n > 1:
                x = line_start_x + int(spacing * i)
            else:
                x = (line_start_x + line_end_x) // 2

            color = event.get("color") or self._default_colors[i % len(self._default_colors)]
            above = (i % 2 == 0)

            # 节点圆点
            r = self._node_radius
            self.add_shape(
                MSO_SHAPE.OVAL,
                x - r, line_y - r, r * 2, r * 2,
                fill_color=color
            )

            # 竖直连接线
            if above:
                conn_start_y = line_y - r - Inches(0.05)
                conn_end_y = line_y - Inches(1.5)
            else:
                conn_start_y = line_y + r + Inches(0.05)
                conn_end_y = line_y + Inches(1.5)

            connector = self.slide.shapes.add_connector(
                1, x, conn_start_y, x, conn_end_y
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(1.5)

            # 事件文本
            label = event.get("label", "")
            date_str = event.get("date", "")
            detail = event.get("detail", "")

            text_x = x - Inches(0.9)
            text_w = Inches(1.8)

            if above:
                text_y = conn_end_y - Inches(0.7)
            else:
                text_y = conn_end_y + Inches(0.05)

            # 日期
            self.add_textbox(
                date_str, text_x, text_y,
                text_w, Inches(0.22),
                font_size=10, bold=True,
                font_color=color,
                alignment=PP_ALIGN.CENTER
            )
            # 标题
            self.add_textbox(
                label, text_x, text_y + Inches(0.22),
                text_w, Inches(0.28),
                font_size=12, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.CENTER
            )
            # 描述
            if detail:
                self.add_textbox(
                    detail, text_x, text_y + Inches(0.5),
                    text_w, Inches(0.35),
                    font_size=9,
                    font_color=RGBColor(127, 140, 141),
                    alignment=PP_ALIGN.CENTER
                )

        # 起止箭头
        arrow_size = Inches(0.15)
        self.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            line_end_x + Inches(0.05), line_y - arrow_size // 2,
            arrow_size, arrow_size,
            fill_color=self._line_color
        )

        return self


if __name__ == "__main__":
    timeline = HorizontalTimeline()
    timeline.set_title("企业年度规划", "Annual Roadmap")
    timeline.set_data([
        {"label": "战略规划", "date": "Q1 2024", "detail": "确定年度目标"},
        {"label": "产品研发", "date": "Q2 2024", "detail": "核心产品开发"},
        {"label": "市场推广", "date": "Q3 2024", "detail": "品牌营销活动"},
        {"label": "成果总结", "date": "Q4 2024", "detail": "年度复盘"},
        {"label": "下年展望", "date": "Q1 2025", "detail": "新周期规划"},
    ])
    timeline.save("output/horizontal_timeline.pptx")
    print("水平时间轴已生成: output/horizontal_timeline.pptx")

"""垂直时间轴 - 垂直线在左侧，事件向右延伸，交替色点"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator


class VerticalTimeline(TemplateGenerator):
    """垂直时间轴"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._events = []  # list of dicts: {label, date, detail, color}
        self._node_radius = Inches(0.12)
        self._line_color = RGBColor(189, 195, 199)
        self._default_colors = [
            RGBColor(52, 152, 219),
            RGBColor(231, 76, 60),
            RGBColor(46, 204, 113),
            RGBColor(155, 89, 182),
            RGBColor(241, 196, 15),
            RGBColor(230, 126, 34),
        ]

    def set_title(self, title, subtitle=""):
        self._title = title
        self._subtitle = subtitle
        return self

    def set_data(self, events):
        """设置事件列表

        Args:
            events: list of dict, 每项包含:
                label (str): 事件名称
                date (str): 日期
                detail (str): 详细描述
                color (RGBColor, optional): 自定义颜色
        """
        self._events = events
        return self

    def generate(self):
        """生成垂直时间轴"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # 标题
        if self._title:
            self.add_textbox(
                self._title, Inches(0.5), Inches(0.2),
                slide_w - Inches(1), Inches(0.5),
                font_size=26, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.LEFT
            )
        if self._subtitle:
            self.add_textbox(
                self._subtitle, Inches(0.5), Inches(0.65),
                slide_w - Inches(1), Inches(0.35),
                font_size=13,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.LEFT
            )

        # 垂直主轴线
        line_x = Inches(2.2)
        top_margin = Inches(1.2)
        bottom_margin = Inches(0.4)
        line_top = int(top_margin)
        line_bottom = int(slide_h - bottom_margin)

        # 主轴线
        axis_line = self.slide.shapes.add_connector(
            1, line_x, line_top, line_x, line_bottom
        )
        axis_line.line.color.rgb = self._line_color
        axis_line.line.width = Pt(3)

        n = len(self._events)
        if n == 0:
            return self

        available_height = line_bottom - line_top
        spacing = available_height / max(n, 1)

        for i, event in enumerate(self._events):
            y = line_top + int(spacing * (i + 0.5))
            color = event.get("color") or self._default_colors[i % len(self._default_colors)]

            # 节点圆点
            r = self._node_radius
            self.add_shape(
                MSO_SHAPE.OVAL,
                line_x - r, y - r, r * 2, r * 2,
                fill_color=color
            )

            # 水平连接线
            connector = self.slide.shapes.add_connector(
                1, line_x + r + Inches(0.05), y,
                line_x + Inches(1.0), y
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(1.5)

            # 日期标签 - 在垂直线左侧
            date_str = event.get("date", "")
            self.add_textbox(
                date_str,
                Inches(0.3), y - Inches(0.15),
                Inches(1.7), Inches(0.3),
                font_size=10, bold=True,
                font_color=color,
                alignment=PP_ALIGN.RIGHT
            )

            # 事件标题和描述 - 在垂直线右侧
            label = event.get("label", "")
            detail = event.get("detail", "")

            text_x = line_x + Inches(1.1)
            text_w = slide_w - text_x - Inches(0.3)

            # 标题
            self.add_textbox(
                label,
                text_x, y - Inches(0.25),
                text_w, Inches(0.3),
                font_size=13, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.LEFT
            )
            # 描述
            if detail:
                self.add_textbox(
                    detail,
                    text_x, y + Inches(0.05),
                    text_w, Inches(0.35),
                    font_size=10,
                    font_color=RGBColor(127, 140, 141),
                    alignment=PP_ALIGN.LEFT
                )

        return self


if __name__ == "__main__":
    timeline = VerticalTimeline()
    timeline.set_title("公司发展历程", "Company History")
    timeline.set_data([
        {"label": "公司成立", "date": "2020-01", "detail": "正式注册成立，核心团队组建"},
        {"label": "天使融资", "date": "2020-06", "detail": "获得天使轮 500 万投资"},
        {"label": "产品上线", "date": "2021-03", "detail": "v1.0 正式发布"},
        {"label": "A轮融资", "date": "2022-01", "detail": "完成 A 轮 3000 万融资"},
        {"label": "团队扩张", "date": "2022-09", "detail": "团队规模突破 100 人"},
        {"label": "国际化", "date": "2023-06", "detail": "进入海外市场"},
        {"label": "B轮融资", "date": "2024-01", "detail": "完成 B 轮融资"},
    ])
    timeline.save("output/vertical_timeline.pptx")
    print("垂直时间轴已生成: output/vertical_timeline.pptx")

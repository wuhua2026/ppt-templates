"""螺旋时间轴 - 节点从中心向外螺旋排列，每个节点为带文字的彩色圆"""

import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from python.base import TemplateGenerator


class SpiralTimeline(TemplateGenerator):
    """螺旋时间轴"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._events = []  # list of dicts: {label, date, detail, color}
        self._node_radius = Inches(0.55)
        self._start_radius = Inches(0.3)
        self._default_colors = [
            RGBColor(52, 152, 219),
            RGBColor(231, 76, 60),
            RGBColor(46, 204, 113),
            RGBColor(155, 89, 182),
            RGBColor(241, 196, 15),
            RGBColor(230, 126, 34),
            RGBColor(26, 188, 156),
            RGBColor(243, 156, 18),
        ]

    def set_title(self, title, subtitle=""):
        self._title = title
        self._subtitle = subtitle
        return self

    def set_data(self, events):
        """设置事件列表

        Args:
            events: list of dict, 每项包含:
                label (str): 事件名称
                date (str): 日期
                detail (str): 详细描述
                color (RGBColor, optional): 自定义颜色
        """
        self._events = events
        return self

    def generate(self):
        """生成螺旋时间轴"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # 标题
        if self._title:
            self.add_textbox(
                self._title, Inches(0.5), Inches(0.2),
                slide_w - Inches(1), Inches(0.5),
                font_size=26, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.LEFT
            )
        if self._subtitle:
            self.add_textbox(
                self._subtitle, Inches(0.5), Inches(0.65),
                slide_w - Inches(1), Inches(0.35),
                font_size=13,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.LEFT
            )

        n = len(self._events)
        if n == 0:
            return self

        # 螺旋中心
        center_x = int(slide_w * 0.48)
        center_y = int(slide_h * 0.55)

        # 螺旋参数
        max_radius = min(slide_w, slide_h) * 0.38
        angle_step = 2 * math.pi / n  # 每个节点间的角度增量
        radius_step = max_radius / max(n, 1)  # 半径递增步长

        prev_x = None
        prev_y = None

        for i, event in enumerate(self._events):
            # 螺旋坐标: 角度递增, 半径递增
            angle = angle_step * i - math.pi / 2  # 从顶部开始
            radius = self._start_radius + radius_step * (i + 1) * 0.5

            x = center_x + int(radius * math.cos(angle))
            y = center_y + int(radius * math.sin(angle))

            color = event.get("color") or self._default_colors[i % len(self._default_colors)]

            # 连接到前一个节点
            if prev_x is not None:
                connector = self.slide.shapes.add_connector(
                    1, prev_x, prev_y, x, y
                )
                connector.line.color.rgb = RGBColor(189, 195, 199)
                connector.line.width = Pt(1.5)
                connector.line.dash_style = 2  # dash

            prev_x = x
            prev_y = y

            # 节点圆形（半径随位置递减）
            r = max(Inches(0.35), self._node_radius - Inches(0.03) * i)
            shape = self.add_shape(
                MSO_SHAPE.OVAL,
                x - r, y - r, r * 2, r * 2,
                fill_color=color
            )

            # 节点内文本 (序号)
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            shape.text_frame.word_wrap = False

            # 标签文本 - 在节点旁
            label = event.get("label", "")
            date_str = event.get("date", "")

            # 根据节点在中心的位置决定文本偏移方向
            dx = x - center_x
            dy = y - center_y
            if abs(dx) > abs(dy):
                # 水平偏移大 - 文本在上下方
                text_offset_x = -Inches(0.7)
                text_offset_y = r + Inches(0.1) if dy >= 0 else -r - Inches(0.45)
            else:
                # 垂直偏移大 - 文本在左右方
                text_offset_x = r + Inches(0.1) if dx >= 0 else -r - Inches(1.5)
                text_offset_y = -Inches(0.15)

            text_x = x + text_offset_x
            text_y = y + text_offset_y

            # 日期
            self.add_textbox(
                date_str, text_x, text_y,
                Inches(1.3), Inches(0.2),
                font_size=9, bold=True,
                font_color=color,
                alignment=PP_ALIGN.CENTER
            )
            # 标签
            self.add_textbox(
                label, text_x, text_y + Inches(0.2),
                Inches(1.3), Inches(0.25),
                font_size=10, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.CENTER
            )

        # 中心标注
        self.add_shape(
            MSO_SHAPE.OVAL,
            center_x - Inches(0.12), center_y - Inches(0.12),
            Inches(0.24), Inches(0.24),
            fill_color=RGBColor(189, 195, 199)
        )

        return self


if __name__ == "__main__":
    timeline = SpiralTimeline()
    timeline.set_title("产品迭代路线图", "Spiral Roadmap")
    timeline.set_data([
        {"label": "概念验证", "date": "2024-01", "detail": "MVP开发"},
        {"label": "Alpha", "date": "2024-03", "detail": "内部测试"},
        {"label": "Beta", "date": "2024-05", "detail": "公测"},
        {"label": "v1.0", "date": "2024-07", "detail": "正式发布"},
        {"label": "v1.5", "date": "2024-09", "detail": "功能扩展"},
        {"label": "v2.0", "date": "2024-11", "detail": "架构升级"},
        {"label": "国际化", "date": "2025-01", "detail": "海外市场"},
        {"label": "生态建设", "date": "2025-03", "detail": "开放平台"},
    ])
    timeline.save("output/spiral_timeline.pptx")
    print("螺旋时间轴已生成: output/spiral_timeline.pptx")

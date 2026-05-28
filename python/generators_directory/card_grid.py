"""卡片网格目录 - 圆角卡片网格，图标区+标题+描述"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, create_grid_positions, add_shadow


class CardGridGenerator(TemplateGenerator):
    """卡片网格目录生成器"""

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self.items = []
        self.title = "目录"
        self.card_colors = [
            hex_to_rgb("#E3F2FD"),
            hex_to_rgb("#E8F5E9"),
            hex_to_rgb("#FFF3E0"),
            hex_to_rgb("#FCE4EC"),
            hex_to_rgb("#F3E5F5"),
            hex_to_rgb("#E0F7FA"),
        ]
        self.accent_colors = [
            hex_to_rgb("#1565C0"),
            hex_to_rgb("#2E7D32"),
            hex_to_rgb("#E65100"),
            hex_to_rgb("#AD1457"),
            hex_to_rgb("#6A1B9A"),
            hex_to_rgb("#00838F"),
        ]

    def set_title(self, title):
        self.title = title

    def set_items(self, items):
        self.items = items

    def generate(self):
        self.create_slide()
        self._add_background()
        self._add_title()
        self._add_cards()
        return self

    def _add_background(self):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5),
            fill_color=hex_to_rgb("#F5F5F5"),
        )

    def _add_title(self):
        self.add_textbox(
            self.title,
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
            font_size=28, font_color=hex_to_rgb("#333333"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _add_cards(self):
        count = len(self.items)
        cols = 3
        rows = (count + cols - 1) // cols
        card_w = Inches(2.6)
        card_h = Inches(2.4)
        gap = Inches(0.4)

        positions = create_grid_positions(
            rows, cols,
            start_x=(Inches(10) - (cols * card_w + (cols - 1) * gap)) // 2,
            start_y=Inches(1.3) + (Inches(5.8) - (rows * card_h + (rows - 1) * gap)) // 2,
            cell_width=card_w,
            cell_height=card_h,
            gap=gap,
        )

        for i, item in enumerate(self.items):
            if i >= len(positions):
                break
            x, y = positions[i]
            color_idx = i % len(self.card_colors)

            # Card background (rounded rectangle)
            card = self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=self.card_colors[color_idx],
            )
            add_shadow(card, blur=Pt(8), offset=Pt(2), opacity=0.15)

            # Icon area (circle at top of card)
            icon_size = Inches(0.7)
            icon_x = x + (card_w - icon_size) // 2
            icon_y = y + Inches(0.25)
            self.add_shape(
                MSO_SHAPE.OVAL,
                icon_x, icon_y, icon_size, icon_size,
                fill_color=self.accent_colors[color_idx],
            )

            # Number inside icon
            txBox = self.slide.shapes.add_textbox(
                icon_x, icon_y + Inches(0.08),
                icon_size, icon_size,
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("number", str(i + 1))
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Card title
            self.add_textbox(
                item.get("title", ""),
                x + Inches(0.15), y + Inches(1.1),
                card_w - Inches(0.3), Inches(0.35),
                font_size=14, font_color=hex_to_rgb("#212121"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

            # Card description
            if item.get("description"):
                self.add_textbox(
                    item["description"],
                    x + Inches(0.15), y + Inches(1.5),
                    card_w - Inches(0.3), Inches(0.6),
                    font_size=10, font_color=hex_to_rgb("#757575"),
                    alignment=PP_ALIGN.CENTER,
                )

            # Bottom accent bar
            bar_w = Inches(0.6)
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + (card_w - bar_w) // 2, y + card_h - Inches(0.25),
                bar_w, Inches(0.04),
                fill_color=self.accent_colors[color_idx],
            )


if __name__ == "__main__":
    gen = CardGridGenerator()
    gen.set_title("内容概览")
    gen.set_items([
        {"number": "1", "title": "行业洞察", "description": "深入了解市场动态与趋势"},
        {"number": "2", "title": "用户需求", "description": "精准把握目标群体痛点"},
        {"number": "3", "title": "竞品策略", "description": "全面分析竞争格局"},
        {"number": "4", "title": "产品方案", "description": "创新解决方案设计"},
        {"number": "5", "title": "运营计划", "description": "高效执行与落地"},
        {"number": "6", "title": "增长路径", "description": "可持续发展策略"},
    ])
    gen.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_card_grid.pptx')
    gen.save(output_path)
    print(f"Saved to {output_path}")

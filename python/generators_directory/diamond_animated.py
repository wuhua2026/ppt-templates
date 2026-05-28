"""菱形动画目录 - 菱形网格排列，交替渐变填充，淡入动画"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, create_grid_positions
from python.animation import apply_animations_to_slide


class DiamondAnimatedGenerator(TemplateGenerator):
    """菱形动画目录生成器"""

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self.items = []
        self.title = "目录"
        self.accent_colors = [
            hex_to_rgb("#2196F3"),
            hex_to_rgb("#1565C0"),
            hex_to_rgb("#0D47A1"),
            hex_to_rgb("#42A5F5"),
        ]

    def set_title(self, title):
        self.title = title

    def set_items(self, items):
        """设置目录项
        items: list of dict with keys 'number', 'title', 'subtitle'
        """
        self.items = items

    def generate(self):
        self.create_slide()
        self._add_background()
        self._add_title()
        diamonds, diamond_ids = self._add_diamonds()
        self._add_decorations()
        self._add_animations(diamond_ids)
        return self

    def _add_background(self):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5),
            fill_color=hex_to_rgb("#F5F7FA"),
        )

    def _add_title(self):
        self.add_textbox(
            self.title,
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
            font_size=32, font_color=hex_to_rgb("#212121"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _add_diamonds(self):
        count = len(self.items)
        cols = 3 if count <= 6 else 4
        rows = (count + cols - 1) // cols
        diamond_size = Inches(1.3)
        gap = Inches(0.6)

        total_w = cols * diamond_size + (cols - 1) * gap
        total_h = rows * diamond_size + (rows - 1) * gap
        start_x = (Inches(10) - total_w) // 2
        start_y = Inches(1.4) + (Inches(5.8) - total_h) // 2

        diamonds = []
        diamond_ids = []

        for i, item in enumerate(self.items):
            row, col = divmod(i, cols)
            cx = start_x + col * (diamond_size + gap) + diamond_size // 2
            cy = start_y + row * (diamond_size + gap) + diamond_size // 2
            left = cx - diamond_size // 2
            top = cy - diamond_size // 2

            shape_id = self._get_next_shape_id()

            # Alternate gradient colors
            color_pair_idx = (i // 2) % len(self.accent_colors)
            c1 = self.accent_colors[color_pair_idx % len(self.accent_colors)]
            c2 = self.accent_colors[(color_pair_idx + 1) % len(self.accent_colors)]

            diamond = self.add_gradient_shape(
                MSO_SHAPE.DIAMOND,
                left, top, diamond_size, diamond_size,
                colors=[c1, c2],
            )

            # Number text inside diamond
            txBox = self.slide.shapes.add_textbox(
                left, top + diamond_size * 0.25,
                diamond_size, diamond_size * 0.35,
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.get("number", str(i + 1))
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Title text below diamond
            self.add_textbox(
                item.get("title", ""),
                cx - Inches(1), cy + diamond_size * 0.45,
                Inches(2), Inches(0.4),
                font_size=12, font_color=hex_to_rgb("#424242"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

            diamonds.append(diamond)
            diamond_ids.append(shape_id)

        return diamonds, diamond_ids

    def _add_decorations(self):
        # Subtle decorative circles in corners
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-0.3), Inches(-0.3),
            Inches(1.2), Inches(1.2),
            fill_color=hex_to_rgb("#E3F2FD"),
        )
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9.3), Inches(6.5),
            Inches(1.2), Inches(1.2),
            fill_color=hex_to_rgb("#E3F2FD"),
        )

    def _add_animations(self, shape_ids):
        animations = []
        for i, sid in enumerate(shape_ids):
            animations.append((sid, "fade_in", i * 200, 600))
        apply_animations_to_slide(self.slide, animations)


if __name__ == "__main__":
    gen = DiamondAnimatedGenerator()
    gen.set_title("菱形动画目录")
    gen.set_items([
        {"number": "01", "title": "项目概述"},
        {"number": "02", "title": "市场分析"},
        {"number": "03", "title": "产品设计"},
        {"number": "04", "title": "技术方案"},
        {"number": "05", "title": "运营策略"},
        {"number": "06", "title": "财务规划"},
    ])
    gen.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_diamond_animated.pptx')
    gen.save(output_path)
    print(f"Saved to {output_path}")

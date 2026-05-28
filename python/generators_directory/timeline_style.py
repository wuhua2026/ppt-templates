"""时间轴目录 - 左侧时间轴线，左右交替分支"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb


class TimelineStyleGenerator(TemplateGenerator):
    """时间轴目录生成器"""

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self.items = []
        self.title = "目录"
        self.colors = [
            hex_to_rgb("#2196F3"),
            hex_to_rgb("#4CAF50"),
            hex_to_rgb("#FF9800"),
            hex_to_rgb("#E91E63"),
            hex_to_rgb("#9C27B0"),
            hex_to_rgb("#00BCD4"),
        ]

    def set_title(self, title):
        self.title = title

    def set_items(self, items):
        self.items = items

    def generate(self):
        self.create_slide()
        self._add_background()
        self._add_title()
        self._add_timeline()
        return self

    def _add_background(self):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5),
            fill_color=hex_to_rgb("#FAFAFA"),
        )

    def _add_title(self):
        self.add_textbox(
            self.title,
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
            font_size=28, font_color=hex_to_rgb("#333333"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _add_timeline(self):
        count = len(self.items)
        if count == 0:
            return

        # Main timeline line position
        line_x = Inches(4.0)
        line_top = Inches(1.2)
        line_bottom = Inches(7.0)
        line_height = line_bottom - line_top
        spacing = line_height / max(count - 1, 1) if count > 1 else 0

        # Draw main vertical line
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_x - Inches(0.02), line_top,
            Inches(0.04), line_height,
            fill_color=hex_to_rgb("#BDBDBD"),
        )

        for i, item in enumerate(self.items):
            y = line_top + i * spacing
            color = self.colors[i % len(self.colors)]
            is_left = (i % 2 == 0)

            # Node circle on timeline
            node_size = Inches(0.4)
            self.add_shape(
                MSO_SHAPE.OVAL,
                line_x - node_size / 2, y - node_size / 2,
                node_size, node_size,
                fill_color=color,
            )
            # Number in node
            txBox = self.slide.shapes.add_textbox(
                line_x - node_size / 2, y - Inches(0.1),
                node_size, node_size,
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("number", str(i + 1))
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Connector line to card
            connector_len = Inches(1.2)
            if is_left:
                conn_start_x = line_x - node_size / 2 - Inches(0.05)
                conn_end_x = conn_start_x - connector_len
            else:
                conn_start_x = line_x + node_size / 2 + Inches(0.05)
                conn_end_x = conn_start_x + connector_len

            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                min(conn_start_x, conn_end_x), y - Inches(0.01),
                connector_len, Inches(0.02),
                fill_color=color,
            )

            # Card / text box
            card_w = Inches(3.5)
            card_h = Inches(0.85)
            if is_left:
                card_x = conn_end_x - card_w + Inches(0.1)
            else:
                card_x = conn_end_x - Inches(0.1)

            card = self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                card_x, y - card_h / 2,
                card_w, card_h,
                fill_color=hex_to_rgb("#FFFFFF"),
                outline_color=hex_to_rgb("#E0E0E0"),
                outline_width=Pt(1),
            )

            # Title in card
            self.add_textbox(
                item.get("title", ""),
                card_x + Inches(0.2), y - card_h / 2 + Inches(0.1),
                card_w - Inches(0.4), Inches(0.3),
                font_size=13, font_color=hex_to_rgb("#212121"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )

            # Subtitle in card
            if item.get("subtitle"):
                self.add_textbox(
                    item["subtitle"],
                    card_x + Inches(0.2), y - card_h / 2 + Inches(0.4),
                    card_w - Inches(0.4), Inches(0.3),
                    font_size=10, font_color=hex_to_rgb("#757575"),
                    alignment=PP_ALIGN.LEFT,
                )

        # Top and bottom endpoints
        for pt_y, pt_color in [(line_top, hex_to_rgb("#2196F3")), (line_bottom, hex_to_rgb("#9E9E9E"))]:
            self.add_shape(
                MSO_SHAPE.OVAL,
                line_x - Inches(0.08), pt_y - Inches(0.08),
                Inches(0.16), Inches(0.16),
                fill_color=pt_color,
            )


if __name__ == "__main__":
    gen = TimelineStyleGenerator()
    gen.set_title("项目里程碑")
    gen.set_items([
        {"number": "01", "title": "需求调研", "subtitle": "第一阶段 - 2周"},
        {"number": "02", "title": "原型设计", "subtitle": "第二阶段 - 3周"},
        {"number": "03", "title": "开发实现", "subtitle": "第三阶段 - 6周"},
        {"number": "04", "title": "测试验收", "subtitle": "第四阶段 - 2周"},
        {"number": "05", "title": "上线发布", "subtitle": "第五阶段 - 1周"},
    ])
    gen.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_timeline_style.pptx')
    gen.save(output_path)
    print(f"Saved to {output_path}")

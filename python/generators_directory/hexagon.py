"""六边形目录 - 六边形网格布局"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb


class HexagonGenerator(TemplateGenerator):
    """六边形目录生成器"""

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self.items = []
        self.title = "目录"
        self.colors = [
            hex_to_rgb("#2196F3"),
            hex_to_rgb("#4CAF50"),
            hex_to_rgb("#FF9800"),
            hex_to_rgb("#E91E63"),
            hex_to_rgb("#9C27B0"),
            hex_to_rgb("#00BCD4"),
        ]

    def set_title(self, title):
        self.title = title

    def set_items(self, items):
        self.items = items

    def generate(self):
        self.create_slide()
        self._add_background()
        self._add_title()
        self._add_hexagons()
        return self

    def _add_background(self):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5),
            fill_color=hex_to_rgb("#F0F4F8"),
        )

    def _add_title(self):
        self.add_textbox(
            self.title,
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
            font_size=28, font_color=hex_to_rgb("#263238"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _add_hexagons(self):
        count = len(self.items)
        hex_size = Inches(1.4)
        h_spacing = hex_size * 0.95
        v_spacing = hex_size * 0.82

        # Calculate grid layout
        cols = 3
        rows = (count + cols - 1) // cols

        total_w = cols * h_spacing
        total_h = rows * v_spacing
        start_x = (Inches(10) - total_w) / 2
        start_y = Inches(1.3) + (Inches(5.8) - total_h) / 2

        for i, item in enumerate(self.items):
            row, col = divmod(i, cols)
            cx = start_x + col * h_spacing + h_spacing / 2
            # Offset even rows
            if row % 2 == 1:
                cx += h_spacing / 2
            cy = start_y + row * v_spacing + v_spacing / 2

            color = self.colors[i % len(self.colors)]
            shape_id = self._get_next_shape_id()

            # Draw hexagon using freeform (approximate with oval + text)
            # python-pptx doesn't have native hexagon, use 6-6 point star (HEXAGON) or approach
            # MSO_SHAPE.HEXAGON exists in python-pptx
            hex_shape = self.add_shape(
                MSO_SHAPE.HEXAGON,
                cx - hex_size / 2, cy - hex_size / 2,
                hex_size, hex_size,
                fill_color=color,
            )

            # Number
            txBox = self.slide.shapes.add_textbox(
                cx - hex_size / 3, cy - Inches(0.25),
                hex_size * 2 / 3, Inches(0.35),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("number", str(i + 1))
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Title below hexagon
            self.add_textbox(
                item.get("title", ""),
                cx - Inches(1), cy + hex_size / 2 + Inches(0.05),
                Inches(2), Inches(0.3),
                font_size=11, font_color=hex_to_rgb("#37474F"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

            # Subtitle
            if item.get("subtitle"):
                self.add_textbox(
                    item["subtitle"],
                    cx - Inches(1), cy + hex_size / 2 + Inches(0.32),
                    Inches(2), Inches(0.25),
                    font_size=9, font_color=hex_to_rgb("#78909C"),
                    alignment=PP_ALIGN.CENTER,
                )


if __name__ == "__main__":
    gen = HexagonGenerator()
    gen.set_title("业务板块")
    gen.set_items([
        {"number": "01", "title": "数据分析", "subtitle": "数据驱动决策"},
        {"number": "02", "title": "用户研究", "subtitle": "深入用户需求"},
        {"number": "03", "title": "产品创新", "subtitle": "持续迭代优化"},
        {"number": "04", "title": "品牌营销", "subtitle": "全渠道推广"},
        {"number": "05", "title": "供应链", "subtitle": "高效协同"},
        {"number": "06", "title": "客户成功", "subtitle": "提升满意度"},
    ])
    gen.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_hexagon.pptx')
    gen.save(output_path)
    print(f"Saved to {output_path}")

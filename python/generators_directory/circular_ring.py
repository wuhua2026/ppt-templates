"""圆环目录 - 中心大圆+外圈标记点+连接线，颜色编码"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, create_circle_positions


class CircularRingGenerator(TemplateGenerator):
    """圆环目录生成器"""

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self.items = []
        self.title = "目录"
        self.colors = [
            hex_to_rgb("#2196F3"),
            hex_to_rgb("#FF9800"),
            hex_to_rgb("#4CAF50"),
            hex_to_rgb("#E91E63"),
            hex_to_rgb("#9C27B0"),
            hex_to_rgb("#00BCD4"),
            hex_to_rgb("#FF5722"),
            hex_to_rgb("#607D8B"),
        ]

    def set_title(self, title):
        self.title = title

    def set_items(self, items):
        self.items = items

    def generate(self):
        self.create_slide()
        self._add_background()
        self._add_title()
        center_x, center_y = Inches(5), Inches(4)
        self._add_center_circle(center_x, center_y)
        self._add_outer_rings_and_connectors(center_x, center_y)
        return self

    def _add_background(self):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5),
            fill_color=hex_to_rgb("#FAFAFA"),
        )

    def _add_title(self):
        self.add_textbox(
            self.title,
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
            font_size=28, font_color=hex_to_rgb("#333333"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _add_center_circle(self, cx, cy):
        radius = Inches(1.2)
        self.add_gradient_shape(
            MSO_SHAPE.OVAL,
            cx - radius, cy - radius,
            radius * 2, radius * 2,
            colors=[hex_to_rgb("#1565C0"), hex_to_rgb("#1976D2")],
        )
        self.add_textbox(
            self.title,
            cx - radius, cy - Inches(0.25),
            radius * 2, Inches(0.5),
            font_size=20, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _add_outer_rings_and_connectors(self, cx, cy):
        count = len(self.items)
        outer_radius = Inches(2.8)
        marker_size = Inches(0.8)
        positions = create_circle_positions(
            cx, cy, outer_radius, count, start_angle=-90
        )

        for i, (item, (px, py)) in enumerate(zip(self.items, positions)):
            color = self.colors[i % len(self.colors)]

            # Connector line from center to marker
            line = self.slide.shapes.add_connector(
                1,  # straight connector
                cx, cy, px, py,
            )
            line.line.color.rgb = color
            line.line.width = Pt(2)
            line.line.dash_style = 2  # dash

            # Outer marker circle
            marker = self.add_shape(
                MSO_SHAPE.OVAL,
                px - marker_size // 2, py - marker_size // 2,
                marker_size, marker_size,
                fill_color=color,
            )

            # Number inside marker
            txBox = self.slide.shapes.add_textbox(
                px - marker_size // 2, py - Inches(0.15),
                marker_size, Inches(0.3),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("number", str(i + 1))
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Title label next to marker
            angle_offset_x = Inches(0.5) if px > cx else Inches(-2.5)
            angle_offset_y = Inches(-0.3) if py < cy else Inches(0.15)
            self.add_textbox(
                item.get("title", ""),
                px + angle_offset_x, py + angle_offset_y,
                Inches(2), Inches(0.35),
                font_size=11, font_color=hex_to_rgb("#333333"),
                bold=True, alignment=PP_ALIGN.LEFT if px > cx else PP_ALIGN.RIGHT,
            )

            # Subtitle below title
            if item.get("subtitle"):
                self.add_textbox(
                    item["subtitle"],
                    px + angle_offset_x, py + angle_offset_y + Inches(0.3),
                    Inches(2), Inches(0.25),
                    font_size=9, font_color=hex_to_rgb("#757575"),
                    alignment=PP_ALIGN.LEFT if px > cx else PP_ALIGN.RIGHT,
                )


if __name__ == "__main__":
    gen = CircularRingGenerator()
    gen.set_title("项目总览")
    gen.set_items([
        {"number": "01", "title": "项目背景", "subtitle": "行业趋势与机遇"},
        {"number": "02", "title": "目标定位", "subtitle": "核心价值主张"},
        {"number": "03", "title": "实施方案", "subtitle": "分阶段推进"},
        {"number": "04", "title": "资源配置", "subtitle": "人力与预算"},
        {"number": "05", "title": "风险管控", "subtitle": "应对策略"},
        {"number": "06", "title": "效果评估", "subtitle": "KPI指标体系"},
    ])
    gen.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_circular_ring.pptx')
    gen.save(output_path)
    print(f"Saved to {output_path}")

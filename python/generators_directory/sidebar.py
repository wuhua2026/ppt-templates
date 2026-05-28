"""侧边栏目录 - 左侧彩色侧边栏 + 右侧章节列表"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb


class SidebarGenerator(TemplateGenerator):
    """侧边栏目录生成器"""

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self.items = []
        self.title = "目录"
        self.accent_color = hex_to_rgb("#1976D2")
        self.sidebar_color = hex_to_rgb("#0D47A1")

    def set_title(self, title):
        self.title = title

    def set_items(self, items):
        self.items = items

    def generate(self):
        self.create_slide()
        self._add_background()
        self._add_sidebar()
        self._add_title_area()
        self._add_section_list()
        return self

    def _add_background(self):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5),
            fill_color=hex_to_rgb("#FFFFFF"),
        )

    def _add_sidebar(self):
        # Left sidebar background
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(2.8), Inches(7.5),
            fill_color=self.sidebar_color,
        )
        # Accent strip
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.8), Inches(0),
            Inches(0.08), Inches(7.5),
            fill_color=hex_to_rgb("#FF6F00"),
        )

    def _add_title_area(self):
        # Title in sidebar
        self.add_textbox(
            self.title,
            Inches(0.3), Inches(0.6),
            Inches(2.2), Inches(0.6),
            font_size=28, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        # Decorative line under title
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3), Inches(1.3),
            Inches(1.0), Inches(0.04),
            fill_color=hex_to_rgb("#FF6F00"),
        )
        # Subtitle
        self.add_textbox(
            "CONTENTS",
            Inches(0.3), Inches(1.5),
            Inches(2.2), Inches(0.4),
            font_size=14, font_color=hex_to_rgb("#90CAF9"),
            alignment=PP_ALIGN.LEFT,
        )

    def _add_section_list(self):
        count = len(self.items)
        list_start_y = Inches(0.8)
        item_height = Inches(1.0)
        max_height = Inches(6.2)
        actual_item_height = min(item_height, max_height / count) if count > 0 else item_height

        for i, item in enumerate(self.items):
            y = list_start_y + i * actual_item_height
            x_num = Inches(3.2)
            x_title = Inches(4.0)

            # Number circle
            circle_size = Inches(0.55)
            color_idx = i % 6
            colors = [
                hex_to_rgb("#2196F3"), hex_to_rgb("#4CAF50"),
                hex_to_rgb("#FF9800"), hex_to_rgb("#E91E63"),
                hex_to_rgb("#9C27B0"), hex_to_rgb("#00BCD4"),
            ]
            circle = self.add_shape(
                MSO_SHAPE.OVAL,
                x_num, y + Inches(0.08),
                circle_size, circle_size,
                fill_color=colors[color_idx],
            )
            # Number in circle
            txBox = self.slide.shapes.add_textbox(
                x_num, y + Inches(0.12),
                circle_size, circle_size,
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("number", f"{i + 1:02d}")
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Section title
            self.add_textbox(
                item.get("title", ""),
                x_title, y + Inches(0.02),
                Inches(5), Inches(0.35),
                font_size=16, font_color=hex_to_rgb("#212121"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )

            # Subtitle
            if item.get("subtitle"):
                self.add_textbox(
                    item["subtitle"],
                    x_title, y + Inches(0.35),
                    Inches(5), Inches(0.3),
                    font_size=11, font_color=hex_to_rgb("#757575"),
                    alignment=PP_ALIGN.LEFT,
                )

            # Separator line
            if i < count - 1:
                self.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x_num + circle_size + Inches(0.2),
                    y + actual_item_height - Inches(0.05),
                    Inches(5.5), Inches(0.01),
                    fill_color=hex_to_rgb("#E0E0E0"),
                )


if __name__ == "__main__":
    gen = SidebarGenerator()
    gen.set_title("目录")
    gen.set_items([
        {"number": "01", "title": "市场概况", "subtitle": "行业规模与增长趋势"},
        {"number": "02", "title": "竞品分析", "subtitle": "主要竞争对手调研"},
        {"number": "03", "title": "用户画像", "subtitle": "目标用户特征"},
        {"number": "04", "title": "产品规划", "subtitle": "功能路线图"},
        {"number": "05", "title": "推广方案", "subtitle": "线上线下策略"},
        {"number": "06", "title": "预算分配", "subtitle": "成本与收益分析"},
    ])
    gen.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_sidebar.pptx')
    gen.save(output_path)
    print(f"Saved to {output_path}")

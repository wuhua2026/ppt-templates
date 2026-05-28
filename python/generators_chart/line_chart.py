"""折线图页面 - 使用连接的圆点和线段绘制折线图，网格线背景"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator


class LineChart(TemplateGenerator):
    """折线图页面"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._series = []  # list of dicts: {name, data: [{x, y}], color}
        self._line_colors = [
            RGBColor(52, 152, 219),
            RGBColor(231, 76, 60),
            RGBColor(46, 204, 113),
            RGBColor(155, 89, 182),
        ]
        self._y_max = None
        self._x_labels = []
        self._y_label = ""
        self._show_markers = True
        self._show_area = False

    def set_title(self, title, subtitle=""):
        self._title = title
        self._subtitle = subtitle
        return self

    def set_data(self, series, x_labels=None, y_label=""):
        """设置折线图数据

        Args:
            series: list of dict, 每项包含:
                name (str): 系列名称
                data (list of float): 数据值列表
                color (RGBColor, optional): 自定义颜色
            x_labels (list of str, optional): X轴标签
            y_label (str): Y轴单位标签
        """
        self._series = series
        self._x_labels = x_labels or []
        self._y_label = y_label
        return self

    def generate(self):
        """生成折线图页面"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # 标题
        if self._title:
            self.add_textbox(
                self._title, Inches(0.5), Inches(0.3),
                slide_w - Inches(1), Inches(0.5),
                font_size=26, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.LEFT
            )
        if self._subtitle:
            self.add_textbox(
                self._subtitle, Inches(0.5), Inches(0.75),
                slide_w - Inches(1), Inches(0.35),
                font_size=13,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.LEFT
            )

        if not self._series:
            return self

        # 图表区域
        chart_left = Inches(1.2)
        chart_top = Inches(1.5)
        chart_right = slide_w - Inches(0.5)
        chart_bottom = slide_h - Inches(1.2)
        chart_width = chart_right - chart_left
        chart_height = chart_bottom - chart_top

        # 计算Y轴最大值
        if self._y_max is None:
            all_vals = []
            for s in self._series:
                all_vals.extend(s["data"])
            self._y_max = max(all_vals) * 1.15 if all_vals else 100

        num_points = max(len(s["data"]) for s in self._series) if self._series else 0

        # 网格线和Y轴刻度
        num_grid_lines = 5
        for i in range(num_grid_lines + 1):
            y_ratio = i / num_grid_lines
            y_pos = chart_bottom - int(chart_height * y_ratio)

            # 水平网格线
            if i > 0:
                grid_line = self.slide.shapes.add_connector(
                    1, chart_left, y_pos, chart_right, y_pos
                )
                grid_line.line.color.rgb = RGBColor(236, 240, 241)
                grid_line.line.width = Pt(0.75)

            # Y轴刻度
            y_val = self._y_max * y_ratio
            self.add_textbox(
                f"{y_val:.0f}",
                chart_left - Inches(0.8), y_pos - Inches(0.12),
                Inches(0.7), Inches(0.24),
                font_size=9,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.RIGHT
            )

        # Y轴标签
        if self._y_label:
            self.add_textbox(
                self._y_label,
                Inches(0.2), chart_top + chart_height // 2 - Inches(0.2),
                Inches(0.9), Inches(0.4),
                font_size=10,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.CENTER
            )

        # X轴基线
        base_line = self.slide.shapes.add_connector(
            1, chart_left, chart_bottom, chart_right, chart_bottom
        )
        base_line.line.color.rgb = RGBColor(189, 195, 199)
        base_line.line.width = Pt(1.5)

        # 垂直网格线
        if num_points > 1:
            for i in range(num_points):
                x = chart_left + int(chart_width * i / (num_points - 1))
                v_line = self.slide.shapes.add_connector(
                    1, x, chart_top, x, chart_bottom
                )
                v_line.line.color.rgb = RGBColor(245, 245, 245)
                v_line.line.width = Pt(0.5)

        # 绘制各系列折线
        for series_idx, series in enumerate(self._series):
            color = series.get("color") or self._line_colors[series_idx % len(self._line_colors)]
            data = series["data"]

            # 数据点坐标
            points = []
            for i, val in enumerate(data):
                if num_points > 1:
                    x = chart_left + int(chart_width * i / (num_points - 1))
                else:
                    x = (chart_left + chart_right) // 2
                y_ratio = val / self._y_max
                y = chart_bottom - int(chart_height * y_ratio)
                points.append((x, y, val))

            # 连接线段
            for i in range(len(points) - 1):
                line = self.slide.shapes.add_connector(
                    1,
                    points[i][0], points[i][1],
                    points[i + 1][0], points[i + 1][1]
                )
                line.line.color.rgb = color
                line.line.width = Pt(2.5)

            # 数据点标记
            if self._show_markers:
                marker_r = Inches(0.08)
                for px, py, val in points:
                    self.add_shape(
                        MSO_SHAPE.OVAL,
                        px - marker_r, py - marker_r,
                        marker_r * 2, marker_r * 2,
                        fill_color=color
                    )

            # 系列名称标签（在第一个点旁）
            if points:
                first_x, first_y = points[0][0], points[0][1]
                self.add_textbox(
                    series.get("name", ""),
                    first_x - Inches(0.5), first_y - Inches(0.3),
                    Inches(1.0), Inches(0.2),
                    font_size=9, bold=True,
                    font_color=color,
                    alignment=PP_ALIGN.CENTER
                )

        # X轴标签
        if self._x_labels and num_points > 0:
            for i, label in enumerate(self._x_labels[:num_points]):
                if num_points > 1:
                    x = chart_left + int(chart_width * i / (num_points - 1))
                else:
                    x = (chart_left + chart_right) // 2
                self.add_textbox(
                    label,
                    x - Inches(0.3), chart_bottom + Inches(0.05),
                    Inches(0.6), Inches(0.25),
                    font_size=9,
                    font_color=RGBColor(85, 85, 85),
                    alignment=PP_ALIGN.CENTER
                )

        # 图例
        if len(self._series) > 1:
            legend_x = int(slide_w * 0.7)
            legend_y = int(chart_bottom + Inches(0.35))
            for si, series in enumerate(self._series):
                color = series.get("color") or self._line_colors[si % len(self._line_colors)]
                lx = legend_x + si * Inches(1.5)
                self.add_shape(
                    MSO_SHAPE.OVAL,
                    lx, legend_y, Inches(0.12), Inches(0.12),
                    fill_color=color
                )
                self.add_textbox(
                    series.get("name", ""),
                    lx + Inches(0.18), legend_y - Inches(0.03),
                    Inches(1.0), Inches(0.2),
                    font_size=9,
                    font_color=RGBColor(85, 85, 85),
                    alignment=PP_ALIGN.LEFT
                )

        return self


if __name__ == "__main__":
    chart = LineChart()
    chart.set_title("季度业绩趋势", "Quarterly Performance Trend")
    chart.set_data(
        series=[
            {"name": "营收", "data": [120, 150, 135, 200, 180, 220, 250]},
            {"name": "成本", "data": [80, 95, 90, 110, 105, 125, 140]},
        ],
        x_labels=["Q1", "Q2", "Q3", "Q4", "Q5", "Q6", "Q7"],
        y_label="万元"
    )
    chart.save("output/line_chart.pptx")
    print("折线图已生成: output/line_chart.pptx")

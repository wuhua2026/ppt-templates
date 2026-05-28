"""动画图表页面 - 带逐条增长动画的柱状图"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.animation import apply_animations_to_slide


class AnimatedBarChart(TemplateGenerator):
    """动画图表页面 - 柱子依次逐条出现"""

    def __init__(self, theme=None):
        super().__init__(theme)
        self._title = ""
        self._subtitle = ""
        self._data = []  # list of dicts: {label, value, color}
        self._bar_colors = [
            RGBColor(52, 152, 219),
            RGBColor(231, 76, 60),
            RGBColor(46, 204, 113),
            RGBColor(155, 89, 182),
            RGBColor(241, 196, 15),
            RGBColor(230, 126, 34),
        ]
        self._y_max = None
        self._y_label = ""
        self._show_values = True
        self._bar_delay = 300  # 每条柱子的延迟 (ms)
        self._bar_duration = 600  # 每条柱子的动画时长 (ms)
        self._animation_type = "grow"  # grow / wipe / fly_in_bottom

    def set_title(self, title, subtitle=""):
        self._title = title
        self._subtitle = subtitle
        return self

    def set_data(self, data, y_label=""):
        """设置图表数据

        Args:
            data: list of dict, 每项包含:
                label (str): X轴标签
                value (float): 数值
                color (RGBColor, optional): 自定义颜色
            y_label (str): Y轴单位标签
        """
        self._data = data
        self._y_label = y_label
        return self

    def set_animation(self, animation_type="grow", delay=300, duration=600):
        """配置动画参数

        Args:
            animation_type (str): 动画类型 grow/wipe/fly_in_bottom
            delay (int): 每条柱子的延迟(ms)
            duration (int): 动画时长(ms)
        """
        self._animation_type = animation_type
        self._bar_delay = delay
        self._bar_duration = duration
        return self

    def generate(self):
        """生成动画图表页面"""
        self.create_slide()
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # 标题
        if self._title:
            self.add_textbox(
                self._title, Inches(0.5), Inches(0.3),
                slide_w - Inches(1), Inches(0.5),
                font_size=26, bold=True,
                font_color=RGBColor(44, 62, 80),
                alignment=PP_ALIGN.LEFT
            )
        if self._subtitle:
            self.add_textbox(
                self._subtitle, Inches(0.5), Inches(0.75),
                slide_w - Inches(1), Inches(0.35),
                font_size=13,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.LEFT
            )

        n = len(self._data)
        if n == 0:
            return self

        # 图表区域
        chart_left = Inches(1.2)
        chart_top = Inches(1.5)
        chart_right = slide_w - Inches(0.5)
        chart_bottom = slide_h - Inches(1.2)
        chart_width = chart_right - chart_left
        chart_height = chart_bottom - chart_top

        # Y轴最大值
        if self._y_max is None:
            max_val = max(d["value"] for d in self._data)
            self._y_max = max_val * 1.15

        # 网格线和Y轴刻度
        num_grid_lines = 5
        for i in range(num_grid_lines + 1):
            y_ratio = i / num_grid_lines
            y_pos = chart_bottom - int(chart_height * y_ratio)

            if i > 0:
                grid_line = self.slide.shapes.add_connector(
                    1, chart_left, y_pos, chart_right, y_pos
                )
                grid_line.line.color.rgb = RGBColor(236, 240, 241)
                grid_line.line.width = Pt(0.75)

            y_val = self._y_max * y_ratio
            self.add_textbox(
                f"{y_val:.0f}",
                chart_left - Inches(0.8), y_pos - Inches(0.12),
                Inches(0.7), Inches(0.24),
                font_size=9,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.RIGHT
            )

        # Y轴标签
        if self._y_label:
            self.add_textbox(
                self._y_label,
                Inches(0.2), chart_top + chart_height // 2 - Inches(0.2),
                Inches(0.9), Inches(0.4),
                font_size=10,
                font_color=RGBColor(127, 140, 141),
                alignment=PP_ALIGN.CENTER
            )

        # X轴基线
        base_line = self.slide.shapes.add_connector(
            1, chart_left, chart_bottom, chart_right, chart_bottom
        )
        base_line.line.color.rgb = RGBColor(189, 195, 199)
        base_line.line.width = Pt(1.5)

        # 绘制柱子并收集动画信息
        animations = []
        bar_gap_ratio = 0.35
        total_gap = chart_width * bar_gap_ratio
        gap = total_gap / (n + 1)
        bar_width = (chart_width - total_gap) / n

        for i, item in enumerate(self._data):
            value = item["value"]
            color = item.get("color") or self._bar_colors[i % len(self._bar_colors)]
            label = item.get("label", "")

            bar_height_ratio = value / self._y_max
            bar_h = int(chart_height * bar_height_ratio)
            bar_x = chart_left + int(gap + i * (bar_width + gap))
            bar_y = chart_bottom - bar_h

            # 渐变柱子
            gradient_end = RGBColor(
                min(255, color[0] + 60),
                min(255, color[1] + 60),
                min(255, color[2] + 60)
            )
            bar_shape = self.add_gradient_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, bar_y, int(bar_width), bar_h,
                [gradient_end, color]
            )

            # 获取柱子形状ID并添加动画
            shape_id = bar_shape.shape_id
            delay = i * self._bar_delay
            animations.append((shape_id, self._animation_type, delay, self._bar_duration))

            # 值标签（也带动画）
            if self._show_values:
                value_box = self.add_textbox(
                    f"{value:.0f}",
                    bar_x, bar_y - Inches(0.25),
                    int(bar_width), Inches(0.22),
                    font_size=10, bold=True,
                    font_color=color,
                    alignment=PP_ALIGN.CENTER
                )
                # 值标签淡入动画
                animations.append((
                    value_box.shape_id, "fade_in",
                    delay + self._bar_duration // 2,
                    300
                ))

            # X轴标签（无动画）
            self.add_textbox(
                label,
                bar_x, chart_bottom + Inches(0.05),
                int(bar_width), Inches(0.3),
                font_size=10,
                font_color=RGBColor(85, 85, 85),
                alignment=PP_ALIGN.CENTER
            )

        # 应用动画到幻灯片
        apply_animations_to_slide(self.slide, animations)

        return self


if __name__ == "__main__":
    chart = AnimatedBarChart()
    chart.set_title("月度业绩增长", "Monthly Growth Animation")
    chart.set_data([
        {"label": "1月", "value": 120},
        {"label": "2月", "value": 95},
        {"label": "3月", "value": 150},
        {"label": "4月", "value": 180},
        {"label": "5月", "value": 210},
        {"label": "6月", "value": 175},
        {"label": "7月", "value": 240},
    ], y_label="万元")
    chart.set_animation(animation_type="grow", delay=300, duration=600)
    chart.save("output/animated_chart.pptx")
    print("动画图表已生成: output/animated_chart.pptx")

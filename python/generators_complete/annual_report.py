"""年度报告完整演示文稿组装器

生成22页年度报告演示文稿：
封面 -> 目录 -> 公司概览 -> 年度成绩(3) -> 数据分析(4) -> 团队建设 -> 社会责任 -> 未来展望 -> 致谢

默认使用红色商务主题，列车穿越迷雾封面。
"""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency


class AnnualReportAssembler(TemplateGenerator):
    """年度报告演示文稿组装器 (22 slides)"""

    DEFAULT_THEME = {
        "primary": "#B71C1C",
        "secondary": "#C62828",
        "accent": "#E53935",
        "highlight": "#EF9A9A",
        "light": "#FFEBEE",
        "bg": "#FAFAFA",
        "text_primary": "#B71C1C",
        "text_secondary": "#455A64",
        "text_light": "#FFFFFF",
    }

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self._title = "年度报告"
        self._subtitle = "Annual Report"
        self._data = {}
        self._theme = theme or self.DEFAULT_THEME

    def set_title(self, title):
        self._title = title
        return self

    def set_subtitle(self, subtitle):
        self._subtitle = subtitle
        return self

    def set_data(self, data):
        """设置演示文稿数据

        data: dict, 可包含以下键:
            - company: 公司信息
            - overview: 公司概览
            - achievements: 年度成绩
            - data_analysis: 数据分析
            - team: 团队建设
            - csr: 社会责任
            - outlook: 未来展望
        """
        self._data = data
        return self

    def _get_color(self, key):
        return hex_to_rgb(getattr(self._theme, key, self._theme.primary))

    # ------------------------------------------------------------------ #
    #  通用辅助方法
    # ------------------------------------------------------------------ #

    def _add_bg(self, color_key="bg"):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            fill_color=self._get_color(color_key),
        )

    def _add_page_title_bar(self, title):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.0),
            fill_color=self._get_color("primary"),
        )
        self.add_textbox(
            title,
            Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7),
            font_size=28, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _add_section_title(self, title, subtitle=""):
        """章节标题页"""
        self.create_slide()
        self._add_bg()
        accent = self._get_color("accent")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Inches(10), Inches(2.0),
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.10)
        self.add_textbox(
            title,
            Inches(1), Inches(3.0), Inches(8), Inches(1.0),
            font_size=36, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self.add_textbox(
                subtitle,
                Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                font_size=18, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _add_thank_you(self):
        self.create_slide()
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [self._get_color("primary"), self._get_color("secondary")],
            angle=135,
        )
        self.add_textbox(
            "Thank You",
            Inches(1), Inches(2.2), Inches(8), Inches(1.2),
            font_size=48, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            "感谢每一位的付出与贡献",
            Inches(1), Inches(3.5), Inches(8), Inches(0.8),
            font_size=22, font_color=RGBColor(239, 154, 154),
            alignment=PP_ALIGN.CENTER,
        )
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.3),
            line_w, Pt(2),
            fill_color=RGBColor(239, 154, 154),
        )
        self.add_textbox(
            "携手共进 再创辉煌",
            Inches(1), Inches(4.2), Inches(8), Inches(0.5),
            font_size=16, font_color=RGBColor(239, 154, 154),
            alignment=PP_ALIGN.CENTER,
        )

    def _add_stat_card(self, left, top, width, height, value, label, color_key="accent"):
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            value,
            left, top + Inches(0.2), width, Inches(0.6),
            font_size=28, font_color=self._get_color(color_key),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            label,
            left, top + Inches(0.85), width, Inches(0.35),
            font_size=12, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.CENTER,
        )

    # ------------------------------------------------------------------ #
    #  各页面生成方法
    # ------------------------------------------------------------------ #

    def _gen_cover(self):
        """封面: 列车穿越迷雾风格"""
        self.create_slide()
        primary = self._get_color("primary")
        secondary = self._get_color("secondary")
        accent = self._get_color("accent")
        highlight = self._get_color("highlight")
        light = self._get_color("light")

        # 背景深色渐变
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [primary, secondary], angle=160,
        )

        # 迷雾效果 - 多层半透明矩形
        fog_layers = [
            (Inches(0), Inches(4.5), Inches(10), Inches(3.0), highlight, 0.85),
            (Inches(0), Inches(5.0), Inches(10), Inches(2.5), highlight, 0.88),
            (Inches(0), Inches(5.5), Inches(10), Inches(2.0), highlight, 0.92),
        ]
        for left, top, w, h, color, trans in fog_layers:
            shape = self.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, w, h,
                fill_color=color,
            )
            set_shape_transparency(shape, trans)

        # 列车轨道线条
        for offset in [Inches(3.5), Inches(4.0), Inches(4.5)]:
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), offset, Inches(10), Pt(2),
                fill_color=RGBColor(255, 255, 255),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.3)

        # 速度线装饰
        for i in range(5):
            line_w = Inches(1.5 - i * 0.2)
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5 + i * 0.8), Inches(3.0 + i * 0.3),
                line_w, Pt(2),
                fill_color=RGBColor(255, 255, 255),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.2 + i * 0.1)

        # 标题
        self.add_textbox(
            self._title,
            Inches(1.5), Inches(1.8), Inches(7.0), Inches(1.2),
            font_size=44, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 年份
        self.add_textbox(
            "2025",
            Inches(1.5), Inches(3.0), Inches(7.0), Inches(0.8),
            font_size=36, font_color=RGBColor(239, 154, 154),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 分隔线
        line_w = Inches(3)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.85),
            line_w, Pt(2),
            fill_color=RGBColor(239, 154, 154),
        )
        # 副标题
        self.add_textbox(
            self._subtitle,
            Inches(1.5), Inches(4.1), Inches(7.0), Inches(0.8),
            font_size=20, font_color=RGBColor(239, 154, 154),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_directory(self):
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目 录")

        sections = [
            ("01", "公司概览"),
            ("02", "年度成绩"),
            ("03", "数据分析"),
            ("04", "团队建设"),
            ("05", "社会责任"),
            ("06", "未来展望"),
        ]
        card_w = Inches(2.6)
        card_h = Inches(1.6)
        gap = Inches(0.4)
        cols = 3
        start_x = (Inches(10) - (cols * card_w + (cols - 1) * gap)) // 2
        start_y = Inches(1.5)

        for i, (num, title) in enumerate(sections):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(0.06),
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                num,
                x, y + Inches(0.2), card_w, Inches(0.5),
                font_size=24, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(0.8), card_w, Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

    def _gen_overview(self):
        """公司概览页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("公司概览")

        overview_items = [
            ("成立时间", "2018年"),
            ("员工规模", "500+"),
            ("业务覆盖", "全国20+城市"),
            ("服务客户", "2000+企业"),
        ]

        box_w = Inches(2.0)
        box_h = Inches(1.5)
        gap = Inches(0.3)
        start_x = (Inches(10) - (4 * box_w + 3 * gap)) // 2

        for i, (label, value) in enumerate(overview_items):
            x = start_x + i * (box_w + gap)
            y = Inches(1.5)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_w, box_h,
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_textbox(
                value,
                x, y + Inches(0.3), box_w, Inches(0.5),
                font_size=24, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                label,
                x, y + Inches(0.9), box_w, Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

        # 公司简介
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(3.3), Inches(9.0), Inches(3.5),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "公司简介",
            Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.4),
            font_size=16, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        desc = (
            "我们是一家专注于企业数字化转型的科技公司，致力于通过人工智能、大数据和云计算技术，"
            "为客户提供一站式数字化解决方案。公司成立7年来，已服务超过2000家企业客户，"
            "涵盖金融、制造、零售、医疗等多个行业。"
        )
        self.add_textbox(
            desc,
            Inches(0.8), Inches(4.1), Inches(8.4), Inches(2.5),
            font_size=13, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.LEFT,
        )

    def _gen_achievement(self, page_num=1):
        """年度成绩页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"年度成绩 - 第{page_num}部分")

        contents = {
            1: ("营收增长", [
                {"value": "2.5亿", "label": "年营收"},
                {"value": "35%", "label": "同比增长"},
                {"value": "85%", "label": "客户续约率"},
            ]),
            2: ("市场拓展", [
                {"value": "20+", "label": "覆盖城市"},
                {"value": "500+", "label": "企业客户"},
                {"value": "50+", "label": "战略伙伴"},
            ]),
            3: ("产品创新", [
                {"value": "12", "label": "新产品上线"},
                {"value": "30+", "label": "技术专利"},
                {"value": "99.9%", "label": "系统可用性"},
            ]),
        }

        title, items = contents.get(page_num, ("年度成绩", []))

        # 大标题
        self.add_textbox(
            title,
            Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.5),
            font_size=20, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # 数据卡片
        box_w = Inches(2.6)
        box_h = Inches(2.0)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * box_w + 2 * gap)) // 2

        for i, item in enumerate(items):
            x = start_x + i * (box_w + gap)
            y = Inches(2.0)
            self._add_stat_card(x, y, box_w, box_h, item["value"], item["label"])

        # 补充说明
        details = {
            1: ["核心业务收入占比 70%", "新兴业务贡献持续提升", "毛利率保持在 60% 以上"],
            2: ["华北、华东、华南三大区域均衡发展", "海外业务开始布局", "渠道合作伙伴持续扩大"],
            3: ["研发投入占营收 20%", "核心技术自主可控", "产品用户满意度 4.8/5.0"],
        }
        detail_items = details.get(page_num, [])
        dy = Inches(4.5)
        for d in detail_items:
            self.add_textbox(
                f"  {d}",
                Inches(1.0), dy, Inches(8.0), Inches(0.35),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            dy += Inches(0.4)

    def _gen_data_analysis(self, page_num=1):
        """数据分析页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"数据分析 - 第{page_num}部分")

        contents = {
            1: ("营收趋势", {
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "values": ["5000万", "6000万", "7000万", "7000万"],
                "total": "2.5亿",
            }),
            2: ("客户分布", {
                "labels": ["金融", "制造", "零售", "医疗", "其他"],
                "values": ["30%", "25%", "20%", "15%", "10%"],
                "total": "2000+客户",
            }),
            3: ("产品使用", {
                "labels": ["数据分析", "自动化", "协作", "API"],
                "values": ["85%", "70%", "60%", "45%"],
                "total": "日活 50万+",
            }),
            4: ("满意度", {
                "labels": ["功能", "性能", "服务", "价格"],
                "values": ["4.8", "4.7", "4.6", "4.5"],
                "total": "综合 4.65/5.0",
            }),
        }

        title, data = contents.get(page_num, ("数据分析", {}))

        # 左侧 - 标题和总览
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.3), Inches(4.2), Inches(5.5),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            title,
            Inches(0.8), Inches(1.5), Inches(3.6), Inches(0.4),
            font_size=18, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        self.add_textbox(
            data.get("total", ""),
            Inches(0.8), Inches(2.1), Inches(3.6), Inches(0.6),
            font_size=28, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        # 模拟图表区域
        chart_h = Inches(3.0)
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(3.0), Inches(3.6), chart_h,
            fill_color=self._get_color("light"),
        )
        set_shape_transparency(self.slide.shapes[-1], 0.3)
        self.add_textbox(
            "图表区域",
            Inches(0.8), Inches(4.2), Inches(3.6), Inches(0.4),
            font_size=12, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.CENTER,
        )

        # 右侧 - 各项数据
        labels = data.get("labels", [])
        values = data.get("values", [])
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.0), Inches(1.3), Inches(4.5), Inches(5.5),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "详细数据",
            Inches(5.3), Inches(1.5), Inches(3.9), Inches(0.4),
            font_size=14, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        for i, (label, value) in enumerate(zip(labels, values)):
            y = Inches(2.2 + i * 1.1)
            self.add_textbox(
                label,
                Inches(5.3), y, Inches(2.0), Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            # 进度条背景
            bar_y = y + Inches(0.35)
            bar_w = Inches(3.5)
            bar_h = Inches(0.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.3), bar_y, bar_w, bar_h,
                fill_color=self._get_color("light"),
            )
            # 进度条前景
            progress = 0.3 + (4 - i) * 0.15
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.3), bar_y, int(bar_w * progress), bar_h,
                fill_color=self._get_color("accent"),
            )
            # 数值
            self.add_textbox(
                value,
                Inches(5.3), bar_y + Inches(0.35), Inches(3.5), Inches(0.25),
                font_size=11, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )

    def _gen_team_building(self):
        """团队建设页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("团队建设")

        stats = [
            {"value": "500+", "label": "团队规模", "growth": "同比增长 25%"},
            {"value": "45%", "label": "硕博占比", "growth": "技术人才为主"},
            {"value": "92%", "label": "员工满意度", "growth": "行业领先水平"},
        ]

        box_w = Inches(2.6)
        box_h = Inches(2.0)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * box_w + 2 * gap)) // 2

        for i, stat in enumerate(stats):
            x = start_x + i * (box_w + gap)
            y = Inches(1.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_w, box_h,
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_textbox(
                stat["value"],
                x, y + Inches(0.2), box_w, Inches(0.5),
                font_size=28, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                stat["label"],
                x, y + Inches(0.75), box_w, Inches(0.3),
                font_size=14, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                stat["growth"],
                x, y + Inches(1.15), box_w, Inches(0.3),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

        # 团队建设举措
        initiatives = [
            ("人才培养", "年度培训计划覆盖全员\n外部专家讲座 20+场"),
            ("文化建设", "季度团建活动\n员工关怀基金"),
            ("激励机制", "股权激励计划\n绩效奖金制度"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(2.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2
        start_y = Inches(3.8)

        for i, (title, desc) in enumerate(initiatives):
            x = start_x + i * (card_w + gap)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, start_y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, start_y, card_w, Inches(0.06),
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                title,
                x, start_y + Inches(0.25), card_w, Inches(0.35),
                font_size=14, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                desc,
                x + Inches(0.3), start_y + Inches(0.8), card_w - Inches(0.6), Inches(1.5),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_csr(self):
        """社会责任页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("社会责任")

        csr_items = [
            ("环境保护", "碳排放减少 30%\n绿色办公全覆盖\n使用清洁能源"),
            ("公益捐赠", "年度公益支出 200万\n教育助学项目\n乡村数字化支持"),
            ("社区参与", "员工志愿者 500+人次\n社区开放日活动\n行业知识分享"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.0)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (title, desc) in enumerate(csr_items):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 图标
            icon_size = Inches(0.8)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - icon_size) // 2, y + Inches(0.3),
                icon_size, icon_size,
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.15)
            self.add_textbox(
                str(i + 1),
                x + (card_w - icon_size) // 2, y + Inches(0.38),
                icon_size, icon_size,
                font_size=24, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(1.3), card_w, Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                desc,
                x + Inches(0.3), y + Inches(1.9), card_w - Inches(0.6), Inches(1.8),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_outlook(self):
        """未来展望页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("未来展望")

        goals = [
            ("2026目标", "营收突破 5 亿\n客户数量翻倍\n海外业务启动"),
            ("战略重点", "AI 产品深耕\n行业解决方案\n生态体系建设"),
            ("组织升级", "团队扩至 800 人\n全球化布局\n组织效能提升"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.0)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (title, desc) in enumerate(goals):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部渐变
            self.add_gradient_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(1.0),
                [self._get_color("primary"), self._get_color("accent")],
                angle=90,
            )
            self.add_textbox(
                title,
                x, y + Inches(0.2), card_w, Inches(0.5),
                font_size=16, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                desc,
                x + Inches(0.3), y + Inches(1.3), card_w - Inches(0.6), Inches(2.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

    # ------------------------------------------------------------------ #
    #  主生成方法
    # ------------------------------------------------------------------ #

    def generate(self):
        """生成完整22页年度报告"""
        self._gen_cover()              # 1. 封面
        self._gen_directory()          # 2. 目录
        self._add_section_title("公司概览", "Company Overview")  # 3. 章节标题
        self._gen_overview()           # 4. 公司概览
        self._add_section_title("年度成绩", "Annual Achievements")  # 5. 章节标题
        self._gen_achievement(1)       # 6. 年度成绩 - 营收增长
        self._gen_achievement(2)       # 7. 年度成绩 - 市场拓展
        self._gen_achievement(3)       # 8. 年度成绩 - 产品创新
        self._add_section_title("业绩亮点", "Performance Highlights")  # 9. 章节标题
        self._add_section_title("数据分析", "Data Analysis")  # 10. 章节标题
        self._gen_data_analysis(1)     # 11. 数据分析 - 营收趋势
        self._gen_data_analysis(2)     # 12. 数据分析 - 客户分布
        self._gen_data_analysis(3)     # 13. 数据分析 - 产品使用
        self._gen_data_analysis(4)     # 14. 数据分析 - 满意度
        self._add_section_title("团队与文化", "Team & Culture")  # 15. 章节标题
        self._gen_team_building()      # 16. 团队建设
        self._add_section_title("社会责任", "CSR")  # 17. 章节标题
        self._gen_csr()                # 18. 社会责任
        self._add_section_title("展望未来", "Future Outlook")  # 19. 章节标题
        self._gen_outlook()            # 20. 未来展望
        self._add_section_title("总结回顾", "Summary")  # 21. 章节标题
        self._add_thank_you()          # 22. 致谢
        return self


if __name__ == "__main__":
    assembler = AnnualReportAssembler()
    assembler.set_title("年度报告").set_subtitle("2025年度公司报告")
    assembler.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_annual_report.pptx')
    assembler.save(output_path)
    print(f"年度报告已生成: {output_path}")

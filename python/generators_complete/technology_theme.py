"""科技主题完整演示文稿组装器

生成16页科技主题演示文稿：
封面 -> 目录 -> 技术架构 -> 核心技术(3) -> 技术栈 -> 性能指标(2) -> 未来规划 -> 致谢

默认使用暗金主题，几何旋转封面(暗色变体)。
"""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency


class TechnologyThemeAssembler(TemplateGenerator):
    """科技主题演示文稿组装器 (16 slides)"""

    DEFAULT_THEME = {
        "primary": "#1A1A2E",
        "secondary": "#16213E",
        "accent": "#D4AF37",
        "highlight": "#F5DEB3",
        "light": "#2C2C3E",
        "bg": "#0F0F1A",
        "text_primary": "#D4AF37",
        "text_secondary": "#A0A0B0",
        "text_light": "#FFFFFF",
    }

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self._title = "科技主题"
        self._subtitle = "Technology Theme"
        self._data = {}
        self._theme = theme or self.DEFAULT_THEME

    def set_title(self, title):
        self._title = title
        return self

    def set_subtitle(self, subtitle):
        self._subtitle = subtitle
        return self

    def set_data(self, data):
        """设置演示文稿数据

        data: dict, 可包含以下键:
            - architecture: 技术架构
            - core_tech: 核心技术列表
            - tech_stack: 技术栈信息
            - performance: 性能指标
            - future_plan: 未来规划
        """
        self._data = data
        return self

    def _get_color(self, key):
        return hex_to_rgb(getattr(self._theme, key, self._theme.primary))

    # ------------------------------------------------------------------ #
    #  通用辅助方法
    # ------------------------------------------------------------------ #

    def _add_bg(self, color_key="bg"):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            fill_color=self._get_color(color_key),
        )

    def _add_page_title_bar(self, title):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.0),
            fill_color=self._get_color("primary"),
        )
        # 金色底边线
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.95), Inches(10), Pt(2),
            fill_color=self._get_color("accent"),
        )
        self.add_textbox(
            title,
            Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7),
            font_size=28, font_color=self._get_color("text_light"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _add_section_title(self, title, subtitle=""):
        """章节标题页"""
        self.create_slide()
        self._add_bg()
        accent = self._get_color("accent")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Inches(10), Inches(2.0),
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.10)
        self.add_textbox(
            title,
            Inches(1), Inches(3.0), Inches(8), Inches(1.0),
            font_size=36, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self.add_textbox(
                subtitle,
                Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                font_size=18, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _add_thank_you(self):
        self.create_slide()
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [self._get_color("primary"), self._get_color("secondary")],
            angle=135,
        )
        self.add_textbox(
            "Thank You",
            Inches(1), Inches(2.2), Inches(8), Inches(1.2),
            font_size=48, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            "探索无界 创新不止",
            Inches(1), Inches(3.5), Inches(8), Inches(0.8),
            font_size=22, font_color=RGBColor(160, 160, 176),
            alignment=PP_ALIGN.CENTER,
        )
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.3),
            line_w, Pt(2),
            fill_color=self._get_color("accent"),
        )
        # 装饰点
        for i in range(5):
            x = Inches(3.5) + Inches(i * 0.8)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x, Inches(4.5), Inches(0.1), Inches(0.1),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.3 + i * 0.12)

    def _add_tech_card(self, left, top, width, height, title, items):
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
            fill_color=self._get_color("light"),
        )
        # 金色顶部线
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, Pt(2),
            fill_color=self._get_color("accent"),
        )
        self.add_textbox(
            title,
            left + Inches(0.2), top + Inches(0.15),
            width - Inches(0.4), Inches(0.35),
            font_size=14, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        y = top + Inches(0.6)
        for item in items[:5]:
            self.add_textbox(
                f"  {item}",
                left + Inches(0.2), y,
                width - Inches(0.4), Inches(0.25),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            y += Inches(0.28)

    # ------------------------------------------------------------------ #
    #  各页面生成方法
    # ------------------------------------------------------------------ #

    def _gen_cover(self):
        """封面: 几何旋转(暗色变体)"""
        self.create_slide()
        primary = self._get_color("primary")
        secondary = self._get_color("secondary")
        accent = self._get_color("accent")
        highlight = self._get_color("highlight")

        # 深色背景
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [primary, secondary], angle=135,
        )

        # 网格线背景效果
        for i in range(8):
            x = Inches(1.25 * i)
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, Inches(0), Pt(1), Inches(7.5),
                fill_color=accent,
            )
            set_shape_transparency(self.slide.shapes[-1], 0.92)
        for i in range(6):
            y = Inches(1.25 * i)
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), y, Inches(10), Pt(1),
                fill_color=accent,
            )
            set_shape_transparency(self.slide.shapes[-1], 0.92)

        # 几何装饰(金色调)
        shapes_config = [
            (MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.2), Inches(2.8), Inches(2.8), 15, accent),
            (MSO_SHAPE.DIAMOND, Inches(5.8), Inches(0.3), Inches(2.2), Inches(2.2), 0, highlight),
            (MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(4.5), Inches(2.5), Inches(2.5), -20, accent),
            (MSO_SHAPE.DIAMOND, Inches(0.8), Inches(5.0), Inches(1.8), Inches(1.8), 30, highlight),
        ]
        for shape_type, left, top, w, h, rot, color in shapes_config:
            shape = self.add_shape(shape_type, left, top, w, h, fill_color=color)
            set_shape_transparency(shape, 0.65)
            if rot != 0:
                from lxml import etree
                xfrm = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rot * 60000)))

        # 中心暗色遮罩
        mask_w = Inches(6.5)
        mask_h = Inches(3.0)
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            (Inches(10) - mask_w) // 2, (Inches(7.5) - mask_h) // 2,
            mask_w, mask_h,
            fill_color=RGBColor(0, 0, 0),
        )
        set_shape_transparency(self.slide.shapes[-1], 0.35)
        self.slide.shapes[-1].line.fill.background()

        # 金色边框装饰
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - mask_w) // 2, (Inches(7.5) - mask_h) // 2,
            mask_w, Pt(1),
            fill_color=accent,
        )
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - mask_w) // 2, (Inches(7.5) + mask_h) // 2 - Pt(1),
            mask_w, Pt(1),
            fill_color=accent,
        )

        # 标题
        self.add_textbox(
            self._title,
            Inches(1.5), Inches(2.6), Inches(7.0), Inches(1.2),
            font_size=44, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 分隔线
        line_w = Inches(2.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.85),
            line_w, Pt(1),
            fill_color=accent,
        )
        # 副标题
        self.add_textbox(
            self._subtitle,
            Inches(1.5), Inches(4.1), Inches(7.0), Inches(0.8),
            font_size=20, font_color=RGBColor(160, 160, 176),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_directory(self):
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目 录")

        sections = [
            ("01", "技术架构"),
            ("02", "核心技术"),
            ("03", "技术栈"),
            ("04", "性能指标"),
            ("05", "未来规划"),
        ]
        card_w = Inches(2.6)
        card_h = Inches(1.6)
        gap = Inches(0.4)
        cols = 3
        start_x = (Inches(10) - (min(cols, len(sections)) * card_w + (min(cols, len(sections)) - 1) * gap)) // 2
        start_y = Inches(1.5)

        for i, (num, title) in enumerate(sections):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=self._get_color("light"),
            )
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Pt(2),
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                num,
                x, y + Inches(0.2), card_w, Inches(0.5),
                font_size=24, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(0.8), card_w, Inches(0.4),
                font_size=14, font_color=self._get_color("text_light"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

    def _gen_architecture(self):
        """技术架构页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("技术架构")

        # 架构层次
        layers = [
            ("应用层", "Web端 | 移动端 | API网关 | 管理后台", self._get_color("accent")),
            ("服务层", "微服务集群 | AI引擎 | 数据处理 | 消息队列", self._get_color("highlight")),
            ("数据层", "关系数据库 | 缓存集群 | 对象存储 | 数据仓库", self._get_color("accent")),
            ("基础设施", "Kubernetes | Docker | CI/CD | 监控告警", self._get_color("highlight")),
        ]

        layer_w = Inches(8.5)
        layer_h = Inches(1.1)
        start_x = (Inches(10) - layer_w) // 2
        start_y = Inches(1.3)

        for i, (title, components, color) in enumerate(layers):
            y = start_y + i * (layer_h + Inches(0.2))
            # 层背景
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                start_x, y, layer_w, layer_h,
                fill_color=self._get_color("light"),
            )
            # 左侧标签
            label_w = Inches(1.5)
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                start_x, y, label_w, layer_h,
                fill_color=color,
            )
            set_shape_transparency(self.slide.shapes[-1], 0.2)
            self.add_textbox(
                title,
                start_x, y + Inches(0.1), label_w, Inches(0.4),
                font_size=14, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 组件
            self.add_textbox(
                components,
                start_x + label_w + Inches(0.3), y + Inches(0.3),
                layer_w - label_w - Inches(0.6), Inches(0.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

        # 连接箭头
        for i in range(3):
            y = start_y + (i + 1) * (layer_h + Inches(0.2)) - Inches(0.15)
            self.add_textbox(
                "|",
                Inches(5), y, Inches(0.5), Inches(0.3),
                font_size=16, font_color=self._get_color("accent"),
                alignment=PP_ALIGN.CENTER,
            )

        # 底部说明
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.3), Inches(9.0), Inches(0.8),
            fill_color=self._get_color("light"),
        )
        self.add_textbox(
            "高可用 | 水平扩展 | 微服务架构 | 云原生 | 安全合规",
            Inches(0.5), Inches(6.45), Inches(9.0), Inches(0.5),
            font_size=13, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    def _gen_core_tech(self, page_num=1):
        """核心技术页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"核心技术 - 第{page_num}部分")

        tech_data = {
            1: [
                {"title": "分布式计算引擎", "items": ["基于Spark/Flink的实时计算", "支持PB级数据处理", "毫秒级延迟保障", "自动容错与恢复"]},
                {"title": "AI推理框架", "items": ["自研模型推理引擎", "支持多模态输入", "GPU集群调度", "模型版本管理"]},
            ],
            2: [
                {"title": "智能路由系统", "items": ["动态负载均衡", "智能流量调度", "故障自动切换", "多活架构支持"]},
                {"title": "安全防护体系", "items": ["零信任安全模型", "端到端加密传输", "实时威胁检测", "合规审计日志"]},
            ],
            3: [
                {"title": "数据湖架构", "items": ["统一数据存储", "Schema自动发现", "数据血缘追踪", "跨源数据融合"]},
                {"title": "边缘计算平台", "items": ["边缘节点管理", "离线计算能力", "云端协同推理", "设备统一接入"]},
            ],
        }

        items = tech_data.get(page_num, tech_data[1])
        card_w = Inches(4.2)
        card_h = Inches(4.5)
        gap = Inches(0.5)
        start_x = (Inches(10) - (2 * card_w + gap)) // 2

        for i, tech in enumerate(items):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            self._add_tech_card(x, y, card_w, card_h, tech["title"], tech["items"])

    def _gen_tech_stack(self):
        """技术栈页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("技术栈")

        stack_layers = [
            ("前端", ["React", "Vue.js", "TypeScript", "Flutter", "WebSocket"]),
            ("后端", ["Go", "Python", "Java", "gRPC", "GraphQL"]),
            ("数据", ["PostgreSQL", "Redis", "Kafka", "Elasticsearch", "ClickHouse"]),
            ("AI/ML", ["PyTorch", "TensorFlow", "HuggingFace", "ONNX", "CUDA"]),
            ("DevOps", ["K8s", "Docker", "Terraform", "ArgoCD", "Prometheus"]),
            ("云服务", ["AWS", "阿里云", "GCP", "Azure", "自建IDC"]),
        ]

        card_w = Inches(2.8)
        card_h = Inches(2.2)
        gap_x = Inches(0.3)
        gap_y = Inches(0.25)
        cols = 3
        start_x = (Inches(10) - (cols * card_w + (cols - 1) * gap_x)) // 2

        for i, (layer_name, techs) in enumerate(stack_layers):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap_x)
            y = Inches(1.3) + row * (card_h + gap_y)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=self._get_color("light"),
            )
            # 层名
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(0.45),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.15)
            self.add_textbox(
                layer_name,
                x, y + Inches(0.05), card_w, Inches(0.35),
                font_size=13, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 技术标签
            tag_y = y + Inches(0.6)
            tag_x = x + Inches(0.15)
            for j, tech in enumerate(techs):
                tag_w = Inches(1.1)
                tag_h = Inches(0.28)
                if tag_x + tag_w > x + card_w:
                    tag_x = x + Inches(0.15)
                    tag_y += Inches(0.35)
                self.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    tag_x, tag_y, tag_w, tag_h,
                    fill_color=self._get_color("primary"),
                )
                self.add_textbox(
                    tech,
                    tag_x, tag_y, tag_w, tag_h,
                    font_size=8, font_color=self._get_color("accent"),
                    bold=True, alignment=PP_ALIGN.CENTER,
                )
                tag_x += tag_w + Inches(0.1)

    def _gen_performance(self, page_num=1):
        """性能指标页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"性能指标 - 第{page_num}部分")

        contents = {
            1: ("系统性能", [
                {"metric": "API响应时间", "value": "<50ms", "pct": "99th"},
                {"metric": "系统可用性", "value": "99.99%", "pct": "SLA"},
                {"metric": "并发处理", "value": "100K+", "pct": "QPS"},
                {"metric": "数据处理", "value": "5PB/天", "pct": "日处理量"},
            ]),
            2: ("AI性能", [
                {"metric": "模型推理", "value": "<10ms", "pct": "P99延迟"},
                {"metric": "训练效率", "value": "3x加速", "pct": "对比基线"},
                {"metric": "准确率", "value": "98.5%", "pct": "F1-Score"},
                {"metric": "模型大小", "value": "50MB", "pct": "压缩后"},
            ]),
        }

        title, metrics = contents.get(page_num, contents[1])

        self.add_textbox(
            title,
            Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.4),
            font_size=18, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        card_w = Inches(2.0)
        card_h = Inches(2.5)
        gap = Inches(0.3)
        start_x = (Inches(10) - (4 * card_w + 3 * gap)) // 2

        for i, m in enumerate(metrics):
            x = start_x + i * (card_w + gap)
            y = Inches(1.8)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=self._get_color("light"),
            )
            # 顶部金色线
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Pt(2),
                fill_color=self._get_color("accent"),
            )
            # 指标名
            self.add_textbox(
                m["metric"],
                x, y + Inches(0.3), card_w, Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )
            # 数值
            self.add_textbox(
                m["value"],
                x, y + Inches(0.8), card_w, Inches(0.6),
                font_size=28, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标签
            self.add_textbox(
                m["pct"],
                x, y + Inches(1.6), card_w, Inches(0.3),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

        # 说明
        details = {
            1: ["全球多区域部署，就近接入", "自动弹性伸缩，应对流量峰值", "全链路监控，秒级告警响应"],
            2: ["自研推理引擎，深度优化", "支持模型量化与剪枝", "在线A/B测试，持续迭代"],
        }
        detail_items = details.get(page_num, [])
        dy = Inches(4.8)
        for d in detail_items:
            self.add_textbox(
                f"  {d}",
                Inches(1.0), dy, Inches(8.0), Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            dy += Inches(0.35)

    def _gen_future_plan(self):
        """未来规划页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("未来规划")

        roadmap = [
            {"time": "2025 H2", "title": "技术深化", "items": ["大模型能力建设", "实时计算升级", "安全体系强化"]},
            {"time": "2026 H1", "title": "平台化", "items": ["开放API平台", "开发者生态", "行业解决方案"]},
            {"time": "2026 H2", "title": "国际化", "items": ["多区域部署", "多语言支持", "合规认证"]},
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.0)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, phase in enumerate(roadmap):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=self._get_color("light"),
            )
            # 顶部金色渐变
            self.add_gradient_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(1.0),
                [self._get_color("accent"), self._get_color("highlight")],
                angle=90,
            )
            set_shape_transparency(self.slide.shapes[-1], 0.2)
            self.add_textbox(
                phase["time"],
                x, y + Inches(0.1), card_w, Inches(0.35),
                font_size=14, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                phase["title"],
                x, y + Inches(0.5), card_w, Inches(0.35),
                font_size=12, font_color=self._get_color("text_light"),
                alignment=PP_ALIGN.CENTER,
            )
            # 列表
            iy = y + Inches(1.3)
            for item in phase["items"]:
                self.add_textbox(
                    f"  {item}",
                    x + Inches(0.2), iy, card_w - Inches(0.4), Inches(0.3),
                    font_size=12, font_color=self._get_color("text_secondary"),
                    alignment=PP_ALIGN.LEFT,
                )
                iy += Inches(0.5)

        # 底部愿景
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(5.8), Inches(9.0), Inches(1.2),
            fill_color=self._get_color("light"),
        )
        self.add_textbox(
            "技术愿景",
            Inches(0.8), Inches(5.9), Inches(8.4), Inches(0.35),
            font_size=14, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        self.add_textbox(
            "以技术创新驱动业务增长，打造行业领先的技术平台，赋能千行百业数字化转型。",
            Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.5),
            font_size=12, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.LEFT,
        )

    # ------------------------------------------------------------------ #
    #  主生成方法
    # ------------------------------------------------------------------ #

    def generate(self):
        """生成完整16页科技主题"""
        self._gen_cover()              # 1. 封面
        self._gen_directory()          # 2. 目录
        self._add_section_title("技术架构", "Architecture")  # 3. 章节标题
        self._gen_architecture()       # 4. 技术架构
        self._add_section_title("核心技术", "Core Technology")  # 5. 章节标题
        self._gen_core_tech(1)         # 6. 核心技术1
        self._gen_core_tech(2)         # 7. 核心技术2
        self._gen_core_tech(3)         # 8. 核心技术3
        self._add_section_title("技术栈与性能", "Tech Stack & Performance")  # 9. 章节标题
        self._gen_tech_stack()         # 10. 技术栈
        self._gen_performance(1)       # 11. 性能指标1
        self._gen_performance(2)       # 12. 性能指标2
        self._add_section_title("未来规划", "Future Roadmap")  # 13. 章节标题
        self._gen_future_plan()        # 14. 未来规划
        self._add_section_title("总结", "Summary")  # 15. 章节标题
        self._add_thank_you()          # 16. 致谢
        return self


if __name__ == "__main__":
    assembler = TechnologyThemeAssembler()
    assembler.set_title("科技主题").set_subtitle("技术架构与未来展望")
    assembler.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_technology_theme.pptx')
    assembler.save(output_path)
    print(f"科技主题已生成: {output_path}")

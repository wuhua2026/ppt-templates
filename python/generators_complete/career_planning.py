"""职业生涯规划完整演示文稿组装器

生成16页职业生涯规划演示文稿：
封面 -> 目录 -> 自我认知(2) -> 职业探索(2) -> 目标设定 -> 能力提升(2) -> 实践经历 -> 成果展示 -> 未来规划 -> 致谢

默认使用紫色渐变主题，圆环封面。
"""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency, create_circle_positions


class CareerPlanningAssembler(TemplateGenerator):
    """职业生涯规划演示文稿组装器 (16 slides)"""

    DEFAULT_THEME = {
        "primary": "#4A148C",
        "secondary": "#6A1B9A",
        "accent": "#9C27B0",
        "highlight": "#CE93D8",
        "light": "#F3E5F5",
        "bg": "#FAFAFA",
        "text_primary": "#4A148C",
        "text_secondary": "#5D4037",
        "text_light": "#FFFFFF",
    }

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self._title = "职业生涯规划"
        self._subtitle = "Career Planning"
        self._data = {}
        self._theme = theme or self.DEFAULT_THEME

    def set_title(self, title):
        self._title = title
        return self

    def set_subtitle(self, subtitle):
        self._subtitle = subtitle
        return self

    def set_data(self, data):
        """设置演示文稿数据

        data: dict, 可包含以下键:
            - self_analysis: 自我分析内容
            - career_exploration: 职业探索内容
            - goals: 目标设定
            - skills: 能力提升计划
            - experiences: 实践经历
            - achievements: 成果展示
            - future_plan: 未来规划
        """
        self._data = data
        return self

    def _get_color(self, key):
        return hex_to_rgb(getattr(self._theme, key, self._theme.primary))

    # ------------------------------------------------------------------ #
    #  通用辅助方法
    # ------------------------------------------------------------------ #

    def _add_bg(self, color_key="bg"):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            fill_color=self._get_color(color_key),
        )

    def _add_page_title_bar(self, title):
        bar_color = self._get_color("primary")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.0),
            fill_color=bar_color,
        )
        self.add_textbox(
            title,
            Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7),
            font_size=28, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _add_section_title(self, title, subtitle=""):
        """章节标题页"""
        self.create_slide()
        self._add_bg()
        accent = self._get_color("accent")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Inches(10), Inches(2.0),
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.10)
        self.add_textbox(
            title,
            Inches(1), Inches(3.0), Inches(8), Inches(1.0),
            font_size=36, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self.add_textbox(
                subtitle,
                Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                font_size=18, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _add_thank_you(self):
        self.create_slide()
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [self._get_color("primary"), self._get_color("secondary")],
            angle=135,
        )
        self.add_textbox(
            "Thank You",
            Inches(1), Inches(2.5), Inches(8), Inches(1.2),
            font_size=48, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            "感谢聆听",
            Inches(1), Inches(3.8), Inches(8), Inches(0.8),
            font_size=24, font_color=RGBColor(206, 147, 216),
            alignment=PP_ALIGN.CENTER,
        )
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.55),
            line_w, Pt(2),
            fill_color=RGBColor(206, 147, 216),
        )

    # ------------------------------------------------------------------ #
    #  各页面生成方法
    # ------------------------------------------------------------------ #

    def _gen_cover(self):
        """封面: 圆环风格"""
        self.create_slide()
        primary = self._get_color("primary")
        secondary = self._get_color("secondary")
        accent = self._get_color("accent")
        highlight = self._get_color("highlight")

        # 背景渐变
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [primary, secondary], angle=135,
        )

        # 中心大圆环
        ring_size = Inches(3.5)
        cx = Inches(5) - ring_size // 2
        cy = Inches(3.75) - ring_size // 2
        self.add_shape(
            MSO_SHAPE.OVAL,
            cx, cy, ring_size, ring_size,
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.7)
        # 内圆
        inner_size = Inches(2.5)
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5) - inner_size // 2, Inches(3.75) - inner_size // 2,
            inner_size, inner_size,
            fill_color=primary,
        )

        # 外圈装饰圆点
        positions = create_circle_positions(Inches(5), Inches(3.75), Inches(3.2), 8, start_angle=0)
        for px, py in positions:
            self.add_shape(
                MSO_SHAPE.OVAL,
                px - Inches(0.15), py - Inches(0.15),
                Inches(0.3), Inches(0.3),
                fill_color=highlight,
            )
            set_shape_transparency(self.slide.shapes[-1], 0.4)

        # 标题
        self.add_textbox(
            self._title,
            Inches(1.5), Inches(2.6), Inches(7.0), Inches(1.2),
            font_size=40, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 分隔线
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.85),
            line_w, Pt(2),
            fill_color=RGBColor(206, 147, 216),
        )
        # 副标题
        self.add_textbox(
            self._subtitle,
            Inches(1.5), Inches(4.1), Inches(7.0), Inches(0.8),
            font_size=20, font_color=RGBColor(206, 147, 216),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_directory(self):
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目 录")

        sections = [
            ("01", "自我认知"),
            ("02", "职业探索"),
            ("03", "目标设定"),
            ("04", "能力提升"),
            ("05", "实践经历"),
            ("06", "成果展示"),
            ("07", "未来规划"),
        ]
        card_w = Inches(2.2)
        card_h = Inches(1.5)
        gap = Inches(0.35)
        cols = 4
        start_x = (Inches(10) - (cols * card_w + (cols - 1) * gap)) // 2
        start_y = Inches(1.6)

        for i, (num, title) in enumerate(sections):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            self._add_content_card(x, y, card_w, card_h, title, num)

    def _add_content_card(self, left, top, width, height, title, num=""):
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, Inches(0.06),
            fill_color=self._get_color("accent"),
        )
        if num:
            circle_size = Inches(0.45)
            self.add_shape(
                MSO_SHAPE.OVAL,
                left + Inches(0.2), top + Inches(0.2),
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                num,
                left + Inches(0.2), top + Inches(0.2),
                circle_size, circle_size,
                font_size=14, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
        title_x = left + Inches(0.75) if num else left + Inches(0.2)
        self.add_textbox(
            title,
            title_x, top + Inches(0.2),
            width - Inches(1.0), Inches(0.4),
            font_size=14, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _gen_self_analysis(self, page_num=1):
        """自我认知页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"自我认知 - 第{page_num}部分")

        contents = {
            1: ("性格特征", [
                "MBTI类型: INTJ / ENTJ",
                "善于逻辑分析与规划",
                "目标导向，执行力强",
                "注重效率与结果",
            ]),
            2: ("兴趣与价值观", [
                "对科技与创新充满热情",
                "追求持续学习与成长",
                "重视团队协作与影响力",
                "渴望创造社会价值",
            ]),
        }
        title, items = contents.get(page_num, ("自我分析", []))

        # 左侧大卡片
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5), Inches(5.0), Inches(5.0),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            title,
            Inches(0.8), Inches(1.7), Inches(4.4), Inches(0.5),
            font_size=20, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        y_pos = Inches(2.4)
        for item in items:
            self.add_textbox(
                f"  {item}",
                Inches(0.8), y_pos, Inches(4.4), Inches(0.35),
                font_size=13, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            y_pos += Inches(0.45)

        # 右侧雷达图占位
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.8), Inches(1.5), Inches(3.7), Inches(5.0),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "能力雷达图",
            Inches(5.8), Inches(1.7), Inches(3.7), Inches(0.4),
            font_size=14, font_color=self._get_color("accent"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 模拟雷达图的圆
        radar_size = Inches(2.5)
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.8) + (Inches(3.7) - radar_size) // 2,
            Inches(2.3) + (Inches(4.0) - radar_size) // 2,
            radar_size, radar_size,
            fill_color=self._get_color("light"),
        )
        set_shape_transparency(self.slide.shapes[-1], 0.5)
        labels = ["领导力", "沟通力", "技术力", "创新力", "执行力"]
        positions = create_circle_positions(
            Inches(7.65), Inches(4.3), Inches(1.2), 5, start_angle=-90
        )
        for (px, py), label in zip(positions, labels):
            self.add_textbox(
                label,
                px - Inches(0.5), py - Inches(0.12),
                Inches(1.0), Inches(0.25),
                font_size=9, font_color=self._get_color("text_secondary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

    def _gen_career_exploration(self, page_num=1):
        """职业探索页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"职业探索 - 第{page_num}部分")

        contents = {
            1: ("行业分析", [
                {"title": "互联网科技", "desc": "持续高增长\n人才需求旺盛"},
                {"title": "人工智能", "desc": "新兴赛道\n前景广阔"},
                {"title": "新能源", "desc": "政策驱动\n长期利好"},
            ]),
            2: ("岗位研究", [
                {"title": "产品经理", "desc": "用户洞察\n需求管理"},
                {"title": "技术总监", "desc": "技术架构\n团队管理"},
                {"title": "创业合伙人", "desc": "商业思维\n资源整合"},
            ]),
        }
        section_title, items = contents.get(page_num, ("职业探索", []))

        card_w = Inches(2.6)
        card_h = Inches(3.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, item in enumerate(items):
            x = start_x + i * (card_w + gap)
            y = Inches(1.5)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 图标区
            icon_size = Inches(0.8)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - icon_size) // 2, y + Inches(0.4),
                icon_size, icon_size,
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.2)
            self.add_textbox(
                f"0{i+1}",
                x + (card_w - icon_size) // 2, y + Inches(0.5),
                icon_size, icon_size,
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标题
            self.add_textbox(
                item["title"],
                x, y + Inches(1.4), card_w, Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                item["desc"],
                x + Inches(0.2), y + Inches(1.9), card_w - Inches(0.4), Inches(1.0),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_goal_setting(self):
        """目标设定页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目标设定")

        goals = [
            ("短期目标\n(1年)", "完成专业技能认证\n积累项目经验\n建立行业人脉"),
            ("中期目标\n(3年)", "成为领域专家\n带领5-10人团队\n年薪突破50万"),
            ("长期目标\n(5-10年)", "行业领军人物\n创办或合伙创业\n实现财务自由"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (title, desc) in enumerate(goals):
            x = start_x + i * (card_w + gap)
            y = Inches(1.5)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部色块
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(1.2),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.1)
            # 编号
            self.add_textbox(
                f"0{i+1}",
                x, y + Inches(0.1), card_w, Inches(0.5),
                font_size=28, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 目标标题
            self.add_textbox(
                title,
                x, y + Inches(0.55), card_w, Inches(0.6),
                font_size=12, font_color=RGBColor(255, 255, 255),
                alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                desc,
                x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), Inches(2.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

    def _gen_skill_improvement(self, page_num=1):
        """能力提升页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"能力提升 - 第{page_num}部分")

        contents = {
            1: ("硬技能提升", [
                {"skill": "数据分析", "level": "Python / SQL / BI工具"},
                {"skill": "项目管理", "level": "PMP认证 / 敏捷方法"},
                {"skill": "技术能力", "level": "AI / 云计算 / 架构设计"},
            ]),
            2: ("软技能提升", [
                {"skill": "领导力", "level": "团队管理与激励"},
                {"skill": "沟通表达", "level": "演讲 / 写作 / 谈判"},
                {"skill": "创新思维", "level": "设计思维 / 创业方法论"},
            ]),
        }
        section_title, items = contents.get(page_num, ("能力提升", []))

        for i, item in enumerate(items):
            y = Inches(1.5) + i * Inches(1.8)
            # 横条卡片
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), y, Inches(9.0), Inches(1.5),
                fill_color=RGBColor(255, 255, 255),
            )
            # 左侧色块
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), y, Inches(0.08), Inches(1.5),
                fill_color=self._get_color("accent"),
            )
            # 技能名称
            self.add_textbox(
                item["skill"],
                Inches(0.9), y + Inches(0.2), Inches(2.5), Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            # 详细
            self.add_textbox(
                item["level"],
                Inches(0.9), y + Inches(0.7), Inches(3.0), Inches(0.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            # 进度条背景
            bar_x = Inches(5.0)
            bar_w = Inches(4.2)
            bar_h = Inches(0.3)
            bar_y = y + Inches(0.6)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_x, bar_y, bar_w, bar_h,
                fill_color=self._get_color("light"),
            )
            # 进度条前景
            progress = 0.5 + i * 0.15
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_x, bar_y, int(bar_w * progress), bar_h,
                fill_color=self._get_color("accent"),
            )

    def _gen_practice(self):
        """实践经历页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("实践经历")

        experiences = [
            ("2023.06 - 2024.06", "XX科技有限公司", "产品实习生", "负责用户增长模块\n参与3个核心功能迭代"),
            ("2022.09 - 2023.05", "校内创业项目", "项目负责人", "带领5人团队完成MVP\n获得校级创业大赛一等奖"),
            ("2021.09 - 2022.06", "学生会", "副主席", "组织10+场大型活动\n管理30人团队"),
        ]

        card_w = Inches(9.0)
        card_h = Inches(1.5)
        start_y = Inches(1.5)

        for i, (time, company, role, desc) in enumerate(experiences):
            y = start_y + i * (card_h + Inches(0.25))
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 左侧色条
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), y, Inches(0.06), card_h,
                fill_color=self._get_color("accent"),
            )
            # 时间
            self.add_textbox(
                time,
                Inches(0.8), y + Inches(0.2), Inches(2.5), Inches(0.3),
                font_size=11, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            # 公司
            self.add_textbox(
                company,
                Inches(0.8), y + Inches(0.5), Inches(2.5), Inches(0.3),
                font_size=14, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            # 职位
            self.add_textbox(
                role,
                Inches(0.8), y + Inches(0.9), Inches(2.5), Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            # 描述
            self.add_textbox(
                desc,
                Inches(4.0), y + Inches(0.3), Inches(5.2), Inches(1.0),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

    def _gen_achievements(self):
        """成果展示页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("成果展示")

        achievements = [
            ("5+", "专业证书"),
            ("3", "创业项目"),
            ("10+", "获奖经历"),
            ("50+", "人脉网络"),
        ]

        box_w = Inches(2.0)
        box_h = Inches(2.2)
        gap = Inches(0.35)
        start_x = (Inches(10) - (4 * box_w + 3 * gap)) // 2

        for i, (value, label) in enumerate(achievements):
            x = start_x + i * (box_w + gap)
            y = Inches(1.8)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_w, box_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部圆
            circle_size = Inches(0.8)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (box_w - circle_size) // 2, y + Inches(0.3),
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                str(i + 1),
                x + (box_w - circle_size) // 2, y + Inches(0.4),
                circle_size, circle_size,
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 数值
            self.add_textbox(
                value,
                x, y + Inches(1.2), box_w, Inches(0.5),
                font_size=28, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标签
            self.add_textbox(
                label,
                x, y + Inches(1.7), box_w, Inches(0.35),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

        # 详情列表
        details = [
            "PMP项目管理专业人士认证",
            "全国大学生创新创业大赛一等奖",
            "Python数据分析高级认证",
            "优秀学生干部 / 奖学金获得者",
        ]
        detail_y = Inches(4.4)
        for detail in details:
            self.add_textbox(
                f"  {detail}",
                Inches(1.0), detail_y, Inches(8.0), Inches(0.35),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            detail_y += Inches(0.4)

    def _gen_future_plan(self):
        """未来规划页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("未来规划")

        phases = [
            ("近期", "1年内", "专业深耕\n技能认证\n项目积累"),
            ("中期", "3年内", "管理转型\n行业影响\n资源积累"),
            ("远期", "5年+", "创业/高管\n行业领军人\n社会贡献"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.0)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (phase, period, desc) in enumerate(phases):
            x = start_x + i * (card_w + gap)
            y = Inches(1.5)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(1.0),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.1)
            self.add_textbox(
                phase,
                x, y + Inches(0.1), card_w, Inches(0.4),
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                period,
                x, y + Inches(0.5), card_w, Inches(0.35),
                font_size=12, font_color=RGBColor(206, 147, 216),
                alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                desc,
                x + Inches(0.3), y + Inches(1.3), card_w - Inches(0.6), Inches(2.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

    # ------------------------------------------------------------------ #
    #  主生成方法
    # ------------------------------------------------------------------ #

    def generate(self):
        """生成完整16页职业生涯规划"""
        self._gen_cover()           # 1. 封面
        self._gen_directory()       # 2. 目录
        self._add_section_title("自我认知", "Self-Awareness")  # 3. 章节标题
        self._gen_self_analysis(1)  # 4. 自我认知1
        self._gen_self_analysis(2)  # 5. 自我认知2
        self._add_section_title("职业探索", "Career Exploration")  # 6. 章节标题
        self._gen_career_exploration(1)  # 7. 职业探索1
        self._gen_career_exploration(2)  # 8. 职业探索2
        self._gen_goal_setting()    # 9. 目标设定
        self._gen_skill_improvement(1)  # 10. 能力提升1
        self._gen_skill_improvement(2)  # 11. 能力提升2
        self._add_section_title("实践与成果", "Practice & Achievements")  # 12. 章节标题
        self._gen_practice()        # 13. 实践经历
        self._gen_achievements()    # 14. 成果展示
        self._gen_future_plan()     # 15. 未来规划
        self._add_thank_you()       # 16. 致谢
        return self


if __name__ == "__main__":
    assembler = CareerPlanningAssembler()
    assembler.set_title("职业生涯规划").set_subtitle("我的职业发展蓝图")
    assembler.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_career_planning.pptx')
    assembler.save(output_path)
    print(f"职业生涯规划已生成: {output_path}")

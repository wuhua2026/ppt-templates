"""商业计划书完整演示文稿组装器

生成18页完整商业计划书演示文稿：
封面 -> 目录 -> 市场分析(3) -> 团队介绍 -> 产品服务(2) -> 商业模式 -> 竞争分析 -> 财务预测(3) -> 融资计划 -> 发展路线 -> 致谢

默认使用蓝色科技主题，几何旋转封面。
"""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency, create_grid_positions


class BusinessPlanAssembler(TemplateGenerator):
    """商业计划书演示文稿组装器 (18 slides)"""

    # 蓝色科技主题默认色
    DEFAULT_THEME = {
        "primary": "#1A237E",
        "secondary": "#283593",
        "accent": "#3F51B5",
        "highlight": "#7986CB",
        "light": "#C5CAE9",
        "bg": "#F5F7FA",
        "text_primary": "#1A237E",
        "text_secondary": "#455A64",
        "text_light": "#FFFFFF",
    }

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self._title = "商业计划书"
        self._subtitle = "Business Plan"
        self._data = {}
        self._theme = theme or self.DEFAULT_THEME

    def set_title(self, title):
        self._title = title
        return self

    def set_subtitle(self, subtitle):
        self._subtitle = subtitle
        return self

    def set_data(self, data):
        """设置演示文稿数据

        data: dict, 可包含以下键:
            - company: 公司名称
            - market_analysis: 市场分析内容列表
            - team_members: 团队成员列表
            - products: 产品服务列表
            - business_model: 商业模式描述
            - competitors: 竞争对手列表
            - financials: 财务预测数据
            - funding: 融资计划
            - roadmap: 发展路线
        """
        self._data = data
        return self

    def _get_color(self, key):
        return hex_to_rgb(getattr(self._theme, key, self._theme.primary))

    # ------------------------------------------------------------------ #
    #  通用页面辅助方法
    # ------------------------------------------------------------------ #

    def _add_bg(self, color_key="bg"):
        """添加纯色背景"""
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            fill_color=self._get_color(color_key),
        )

    def _add_page_title_bar(self, title):
        """页面顶部标题栏"""
        bar_color = self._get_color("primary")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.0),
            fill_color=bar_color,
        )
        self.add_textbox(
            title,
            Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7),
            font_size=28, font_color=self._get_color("text_light"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _add_section_title(self, title, subtitle=""):
        """章节标题页（居中大标题+副标题）"""
        self.create_slide()
        self._add_bg()
        accent = self._get_color("accent")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Inches(10), Inches(2.0),
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.10)
        self.add_textbox(
            title,
            Inches(1), Inches(3.0), Inches(8), Inches(1.0),
            font_size=36, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self.add_textbox(
                subtitle,
                Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                font_size=18, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _add_content_card(self, left, top, width, height, title, items, icon_num=""):
        """内容卡片"""
        card_color = RGBColor(255, 255, 255)
        card = self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
            fill_color=card_color,
        )
        set_shape_transparency(card, 0.05)
        # 顶部色条
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, Inches(0.06),
            fill_color=self._get_color("accent"),
        )
        # 编号圆圈
        if icon_num:
            circle_size = Inches(0.45)
            self.add_shape(
                MSO_SHAPE.OVAL,
                left + Inches(0.2), top + Inches(0.2),
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                icon_num,
                left + Inches(0.2), top + Inches(0.2),
                circle_size, circle_size,
                font_size=14, font_color=self._get_color("text_light"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
        # 标题
        title_x = left + Inches(0.75) if icon_num else left + Inches(0.2)
        self.add_textbox(
            title,
            title_x, top + Inches(0.2),
            width - Inches(1.0), Inches(0.4),
            font_size=14, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        # 列表项
        item_top = top + Inches(0.7)
        for item in items[:4]:
            self.add_textbox(
                f"  {item}",
                left + Inches(0.25), item_top,
                width - Inches(0.5), Inches(0.3),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            item_top += Inches(0.28)

    def _add_thank_you(self):
        """致谢页"""
        self.create_slide()
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [self._get_color("primary"), self._get_color("secondary")],
            angle=135,
        )
        self.add_textbox(
            "Thank You",
            Inches(1), Inches(2.5), Inches(8), Inches(1.2),
            font_size=48, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            "感谢您的聆听",
            Inches(1), Inches(3.8), Inches(8), Inches(0.8),
            font_size=24, font_color=RGBColor(197, 202, 233),
            alignment=PP_ALIGN.CENTER,
        )
        # 装饰线
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.55),
            line_w, Pt(2),
            fill_color=RGBColor(197, 202, 233),
        )

    # ------------------------------------------------------------------ #
    #  各页面生成方法
    # ------------------------------------------------------------------ #

    def _gen_cover(self):
        """封面: 几何旋转风格"""
        self.create_slide()
        primary = self._get_color("primary")
        secondary = self._get_color("secondary")
        accent = self._get_color("accent")
        highlight = self._get_color("highlight")
        light = self._get_color("light")

        # 背景渐变
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [primary, secondary], angle=45,
        )

        # 几何装饰
        shapes_config = [
            (MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(3.0), Inches(3.0), 15, accent),
            (MSO_SHAPE.DIAMOND, Inches(5.5), Inches(0.1), Inches(2.5), Inches(2.5), 0, highlight),
            (MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(4.0), Inches(2.8), Inches(2.8), -20, light),
            (MSO_SHAPE.DIAMOND, Inches(1.0), Inches(4.5), Inches(2.2), Inches(2.2), 30, accent),
            (MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.0), Inches(2.5), Inches(2.5), 45, highlight),
        ]
        for shape_type, left, top, w, h, rot, color in shapes_config:
            shape = self.add_shape(shape_type, left, top, w, h, fill_color=color)
            set_shape_transparency(shape, 0.60)
            if rot != 0:
                from lxml import etree
                xfrm = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rot * 60000)))

        # 标题
        self.add_textbox(
            self._title,
            Inches(1.5), Inches(2.6), Inches(7.0), Inches(1.2),
            font_size=42, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 分隔线
        line_w = Inches(2.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.85),
            line_w, Pt(2),
            fill_color=RGBColor(197, 202, 233),
        )
        # 副标题
        self.add_textbox(
            self._subtitle,
            Inches(1.5), Inches(4.1), Inches(7.0), Inches(0.8),
            font_size=20, font_color=RGBColor(197, 202, 233),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_directory(self):
        """目录页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目 录")

        sections = [
            ("01", "市场分析"),
            ("02", "团队介绍"),
            ("03", "产品服务"),
            ("04", "商业模式"),
            ("05", "竞争分析"),
            ("06", "财务预测"),
            ("07", "融资计划"),
            ("08", "发展路线"),
        ]
        card_w = Inches(2.0)
        card_h = Inches(1.4)
        gap = Inches(0.35)
        cols = 4
        start_x = (Inches(10) - (cols * card_w + (cols - 1) * gap)) // 2
        start_y = Inches(1.6)

        for i, (num, title) in enumerate(sections):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            self._add_content_card(x, y, card_w, card_h, title, [], icon_num=num)

    def _gen_market_analysis(self, page_num=1):
        """市场分析页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"市场分析 - 第{page_num}部分")

        contents = {
            1: ("行业概览", ["全球市场规模持续增长", "年复合增长率 15%+", "数字化转型加速推进", "政策环境利好"]),
            2: ("目标市场", ["一线及新一线城市", "25-45岁核心用户群", "B端企业客户为主", "C端高端消费者"]),
            3: ("用户需求", ["效率提升需求强烈", "个性化定制趋势", "一站式解决方案", "数据安全与隐私"]),
        }
        title, items = contents.get(page_num, ("市场分析", []))
        self._add_content_card(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5), title, items)

        # 数据亮点
        highlights = [
            ("100亿+", "市场规模"),
            ("15%", "年增长率"),
            ("5000万", "目标用户"),
        ]
        box_w = Inches(2.5)
        start_x = (Inches(10) - (3 * box_w + 2 * Inches(0.5))) // 2
        for i, (value, label) in enumerate(highlights):
            x = start_x + i * (box_w + Inches(0.5))
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, Inches(5.3), box_w, Inches(1.5),
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_textbox(
                value,
                x, Inches(5.4), box_w, Inches(0.8),
                font_size=28, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                label,
                x, Inches(6.1), box_w, Inches(0.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_team(self):
        """团队介绍页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("团队介绍")

        team = self._data.get("team_members", [
            {"name": "张三", "role": "CEO", "desc": "10年行业经验"},
            {"name": "李四", "role": "CTO", "desc": "技术架构专家"},
            {"name": "王五", "role": "CMO", "desc": "品牌营销资深"},
            {"name": "赵六", "role": "CFO", "desc": "金融投资背景"},
        ])

        card_w = Inches(2.0)
        card_h = Inches(3.2)
        gap = Inches(0.3)
        start_x = (Inches(10) - (len(team[:4]) * card_w + (len(team[:4]) - 1) * gap)) // 2

        for i, member in enumerate(team[:4]):
            x = start_x + i * (card_w + gap)
            y = Inches(1.5)
            # 卡片背景
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 头像占位圆
            circle_size = Inches(1.0)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - circle_size) // 2, y + Inches(0.3),
                circle_size, circle_size,
                fill_color=self._get_color("light"),
            )
            self.add_textbox(
                member["name"][0],
                x + (card_w - circle_size) // 2, y + Inches(0.45),
                circle_size, circle_size,
                font_size=28, font_color=self._get_color("primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 姓名
            self.add_textbox(
                member["name"],
                x, y + Inches(1.5), card_w, Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 职位
            self.add_textbox(
                member["role"],
                x, y + Inches(1.9), card_w, Inches(0.35),
                font_size=12, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                member.get("desc", ""),
                x, y + Inches(2.3), card_w, Inches(0.5),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_product(self, page_num=1):
        """产品服务页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"产品服务 - 第{page_num}部分")

        products = {
            1: ("核心产品", [
                {"title": "智能分析平台", "desc": "AI驱动的数据分析"},
                {"title": "云服务解决方案", "desc": "弹性可扩展架构"},
                {"title": "移动应用套件", "desc": "全场景覆盖"},
            ]),
            2: ("服务体系", [
                {"title": "技术咨询", "desc": "定制化技术方案"},
                {"title": "运维支持", "desc": "7x24小时保障"},
                {"title": "培训赋能", "desc": "知识转移与培养"},
            ]),
        }
        section_title, items = products.get(page_num, ("产品服务", []))

        card_w = Inches(2.6)
        card_h = Inches(3.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, item in enumerate(items):
            x = start_x + i * (card_w + gap)
            y = Inches(1.5)
            # 卡片
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部色块
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(0.8),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.15)
            # 图标圆
            icon_size = Inches(0.7)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - icon_size) // 2, y + Inches(0.05),
                icon_size, icon_size,
                fill_color=RGBColor(255, 255, 255),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.3)
            self.add_textbox(
                str(i + 1),
                x + (card_w - icon_size) // 2, y + Inches(0.12),
                icon_size, icon_size,
                font_size=20, font_color=self._get_color("primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标题
            self.add_textbox(
                item["title"],
                x, y + Inches(1.1), card_w, Inches(0.5),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                item["desc"],
                x + Inches(0.2), y + Inches(1.7), card_w - Inches(0.4), Inches(1.2),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

        self.add_textbox(
            section_title,
            Inches(0.6), Inches(5.8), Inches(8.8), Inches(0.5),
            font_size=14, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_business_model(self):
        """商业模式页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("商业模式")

        models = [
            ("SaaS订阅", "按月/年收费\n灵活定价方案"),
            ("增值服务", "高级功能\n定制开发"),
            ("数据服务", "行业洞察\n数据分析"),
            ("生态合作", "渠道分成\n联合推广"),
        ]

        card_w = Inches(2.0)
        card_h = Inches(2.8)
        gap = Inches(0.3)
        start_x = (Inches(10) - (4 * card_w + 3 * gap)) // 2

        for i, (title, desc) in enumerate(models):
            x = start_x + i * (card_w + gap)
            y = Inches(1.6)
            # 卡片
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部圆
            circle_size = Inches(0.9)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - circle_size) // 2, y + Inches(0.3),
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                f"0{i+1}",
                x + (card_w - circle_size) // 2, y + Inches(0.45),
                circle_size, circle_size,
                font_size=22, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标题
            self.add_textbox(
                title,
                x, y + Inches(1.4), card_w, Inches(0.4),
                font_size=14, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                desc,
                x + Inches(0.15), y + Inches(1.8), card_w - Inches(0.3), Inches(0.8),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_competition(self):
        """竞争分析页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("竞争分析")

        # 对比表头
        cols = ["维度", "我方", "竞品A", "竞品B"]
        col_w = [Inches(1.8), Inches(2.4), Inches(2.4), Inches(2.4)]
        start_x = Inches(0.5)
        header_y = Inches(1.5)

        for i, (col_title, w) in enumerate(zip(cols, col_w)):
            x = start_x + sum(c.inches for c in [Inches(0)] + col_w[:i])
            # 修正计算
            x_pos = start_x
            for j in range(i):
                x_pos += col_w[j]
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x_pos, header_y, w, Inches(0.5),
                fill_color=self._get_color("primary") if i == 0 else self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.1 if i == 0 else 0.0)
            self.add_textbox(
                col_title,
                x_pos, header_y, w, Inches(0.5),
                font_size=13, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

        rows = [
            ("技术实力", "领先", "中等", "中等"),
            ("市场份额", "快速增长", "稳定", "下降"),
            ("用户体验", "优秀", "良好", "一般"),
            ("价格策略", "灵活", "偏高", "低价"),
            ("创新能力", "强", "一般", "弱"),
        ]

        for r, (dim, ours, a, b) in enumerate(rows):
            y = header_y + Inches(0.5 + r * 0.55)
            row_bg = RGBColor(245, 247, 250) if r % 2 == 0 else RGBColor(255, 255, 255)
            values = [dim, ours, a, b]
            for i, (val, w) in enumerate(zip(values, col_w)):
                x_pos = start_x
                for j in range(i):
                    x_pos += col_w[j]
                self.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x_pos, y, w, Inches(0.5),
                    fill_color=row_bg,
                )
                fc = self._get_color("text_primary") if i == 0 else self._get_color("text_secondary")
                self.add_textbox(
                    val,
                    x_pos, y, w, Inches(0.5),
                    font_size=11, font_color=fc,
                    bold=(i == 0), alignment=PP_ALIGN.CENTER,
                )

    def _gen_financial(self, page_num=1):
        """财务预测页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"财务预测 - 第{page_num}部分")

        contents = {
            1: ("收入预测", [
                {"year": "2025", "revenue": "500万", "growth": "-"},
                {"year": "2026", "revenue": "2000万", "growth": "300%"},
                {"year": "2027", "revenue": "8000万", "growth": "300%"},
            ]),
            2: ("成本结构", [
                {"item": "研发投入", "ratio": "40%", "amount": "800万"},
                {"item": "市场推广", "ratio": "25%", "amount": "500万"},
                {"item": "运营成本", "ratio": "20%", "amount": "400万"},
            ]),
            3: ("盈利预测", [
                {"year": "2025", "profit": "-300万", "rate": "-60%"},
                {"year": "2026", "profit": "200万", "rate": "10%"},
                {"year": "2027", "profit": "3000万", "rate": "37.5%"},
            ]),
        }

        title, items = contents.get(page_num, ("财务数据", []))
        # 标题
        self.add_textbox(
            title,
            Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.5),
            font_size=18, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        # 数据卡片
        card_w = Inches(2.6)
        card_h = Inches(3.8)
        gap = Inches(0.4)
        start_x = (Inches(10) - (len(items) * card_w + (len(items) - 1) * gap)) // 2

        for i, item in enumerate(items):
            x = start_x + i * (card_w + gap)
            y = Inches(2.0)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部色条
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(0.08),
                fill_color=self._get_color("accent"),
            )
            # 内容
            values = list(item.values())
            for j, val in enumerate(values):
                fc = self._get_color("text_primary") if j == 0 else self._get_color("text_secondary")
                fs = 14 if j == 0 else 22 if j == 1 else 12
                bld = j <= 1
                self.add_textbox(
                    val,
                    x + Inches(0.2), y + Inches(0.4 + j * 1.0),
                    card_w - Inches(0.4), Inches(0.6),
                    font_size=fs, font_color=fc,
                    bold=bld, alignment=PP_ALIGN.CENTER,
                )

    def _gen_funding(self):
        """融资计划页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("融资计划")

        rounds = [
            {"round": "天使轮", "amount": "500万", "use": "产品研发", "time": "2025 Q1"},
            {"round": "Pre-A轮", "amount": "2000万", "use": "市场拓展", "time": "2025 Q4"},
            {"round": "A轮", "amount": "8000万", "use": "规模化", "time": "2026 Q3"},
        ]

        card_w = Inches(2.6)
        card_h = Inches(3.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, rd in enumerate(rounds):
            x = start_x + i * (card_w + gap)
            y = Inches(1.6)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 金额圆
            circle_size = Inches(1.2)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - circle_size) // 2, y + Inches(0.3),
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                rd["amount"],
                x + (card_w - circle_size) // 2, y + Inches(0.55),
                circle_size, circle_size,
                font_size=14, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 轮次
            self.add_textbox(
                rd["round"],
                x, y + Inches(1.65), card_w, Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 用途
            self.add_textbox(
                f"用途: {rd['use']}",
                x, y + Inches(2.1), card_w, Inches(0.35),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )
            # 时间
            self.add_textbox(
                f"时间: {rd['time']}",
                x, y + Inches(2.5), card_w, Inches(0.35),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_roadmap(self):
        """发展路线页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("发展路线")

        phases = [
            ("2025 Q1", "产品立项", "完成核心团队组建\n启动产品研发"),
            ("2025 Q3", "产品上线", "MVP发布\n获取种子用户"),
            ("2026 Q1", "市场拓展", "规模化推广\n建立渠道合作"),
            ("2026 Q3", "行业领先", "市场份额前三\n启动B轮融资"),
        ]

        # 时间轴线
        line_y = Inches(3.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), line_y, Inches(8.4), Pt(3),
            fill_color=self._get_color("accent"),
        )

        node_w = Inches(2.0)
        total_w = len(phases) * node_w + (len(phases) - 1) * Inches(0.3)
        start_x = (Inches(10) - total_w) // 2

        for i, (time, title, desc) in enumerate(phases):
            x = start_x + i * (node_w + Inches(0.3))
            # 节点圆
            circle_size = Inches(0.5)
            cx = x + node_w // 2 - circle_size // 2
            self.add_shape(
                MSO_SHAPE.OVAL,
                cx, line_y - circle_size // 2,
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                str(i + 1),
                cx, line_y - Inches(0.15),
                circle_size, circle_size,
                font_size=14, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

            # 上方卡片（偶数索引）
            if i % 2 == 0:
                card_y = line_y - Inches(2.2)
                self.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, card_y, node_w, Inches(1.8),
                    fill_color=RGBColor(255, 255, 255),
                )
                self.add_textbox(
                    time,
                    x, card_y + Inches(0.1), node_w, Inches(0.3),
                    font_size=11, font_color=self._get_color("accent"),
                    bold=True, alignment=PP_ALIGN.CENTER,
                )
                self.add_textbox(
                    title,
                    x, card_y + Inches(0.4), node_w, Inches(0.35),
                    font_size=13, font_color=self._get_color("text_primary"),
                    bold=True, alignment=PP_ALIGN.CENTER,
                )
                self.add_textbox(
                    desc,
                    x + Inches(0.1), card_y + Inches(0.8), node_w - Inches(0.2), Inches(0.8),
                    font_size=9, font_color=self._get_color("text_secondary"),
                    alignment=PP_ALIGN.CENTER,
                )
            else:
                # 下方卡片（奇数索引）
                card_y = line_y + Inches(0.6)
                self.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, card_y, node_w, Inches(1.8),
                    fill_color=RGBColor(255, 255, 255),
                )
                self.add_textbox(
                    time,
                    x, card_y + Inches(0.1), node_w, Inches(0.3),
                    font_size=11, font_color=self._get_color("accent"),
                    bold=True, alignment=PP_ALIGN.CENTER,
                )
                self.add_textbox(
                    title,
                    x, card_y + Inches(0.4), node_w, Inches(0.35),
                    font_size=13, font_color=self._get_color("text_primary"),
                    bold=True, alignment=PP_ALIGN.CENTER,
                )
                self.add_textbox(
                    desc,
                    x + Inches(0.1), card_y + Inches(0.8), node_w - Inches(0.2), Inches(0.8),
                    font_size=9, font_color=self._get_color("text_secondary"),
                    alignment=PP_ALIGN.CENTER,
                )

    # ------------------------------------------------------------------ #
    #  主生成方法
    # ------------------------------------------------------------------ #

    def generate(self):
        """生成完整18页商业计划书"""
        self._gen_cover()              # 1. 封面
        self._gen_directory()          # 2. 目录
        self._add_section_title("市场分析", "Market Analysis")  # 3. 章节标题
        self._gen_market_analysis(1)   # 4. 市场分析1
        self._gen_market_analysis(2)   # 5. 市场分析2
        self._gen_market_analysis(3)   # 6. 市场分析3
        self._gen_team()               # 7. 团队介绍
        self._gen_product(1)           # 8. 产品服务1
        self._gen_product(2)           # 9. 产品服务2
        self._gen_business_model()     # 10. 商业模式
        self._gen_competition()        # 11. 竞争分析
        self._add_section_title("财务预测", "Financial Forecast")  # 12. 章节标题
        self._gen_financial(1)         # 13. 财务预测1
        self._gen_financial(2)         # 14. 财务预测2
        self._gen_financial(3)         # 15. 财务预测3
        self._gen_funding()            # 16. 融资计划
        self._gen_roadmap()            # 17. 发展路线
        self._add_thank_you()          # 18. 致谢
        return self


if __name__ == "__main__":
    assembler = BusinessPlanAssembler()
    assembler.set_title("商业计划书").set_subtitle("XX科技有限公司 2025年度")
    assembler.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_business_plan.pptx')
    assembler.save(output_path)
    print(f"商业计划书已生成: {output_path}")

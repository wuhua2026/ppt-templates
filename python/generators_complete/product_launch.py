"""产品发布会完整演示文稿组装器

生成20页产品发布会演示文稿：
封面 -> 目录 -> 痛点分析 -> 解决方案 -> 产品功能(4) -> 产品演示 -> 用户案例(2) -> 价格方案 -> 竞品对比 -> 发布计划 -> 致谢

默认使用海洋蓝主题，菱形展示封面。
"""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency


class ProductLaunchAssembler(TemplateGenerator):
    """产品发布会演示文稿组装器 (20 slides)"""

    DEFAULT_THEME = {
        "primary": "#0D47A1",
        "secondary": "#1565C0",
        "accent": "#2196F3",
        "highlight": "#64B5F6",
        "light": "#BBDEFB",
        "bg": "#F0F7FF",
        "text_primary": "#0D47A1",
        "text_secondary": "#455A64",
        "text_light": "#FFFFFF",
    }

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self._title = "产品发布会"
        self._subtitle = "Product Launch"
        self._data = {}
        self._theme = theme or self.DEFAULT_THEME

    def set_title(self, title):
        self._title = title
        return self

    def set_subtitle(self, subtitle):
        self._subtitle = subtitle
        return self

    def set_data(self, data):
        """设置演示文稿数据

        data: dict, 可包含以下键:
            - pain_points: 痛点列表
            - solution: 解决方案
            - features: 产品功能列表
            - demo_url: 演示链接
            - cases: 用户案例
            - pricing: 价格方案
            - competitors: 竞品信息
            - launch_plan: 发布计划
        """
        self._data = data
        return self

    def _get_color(self, key):
        return hex_to_rgb(getattr(self._theme, key, self._theme.primary))

    # ------------------------------------------------------------------ #
    #  通用辅助方法
    # ------------------------------------------------------------------ #

    def _add_bg(self, color_key="bg"):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            fill_color=self._get_color(color_key),
        )

    def _add_page_title_bar(self, title):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.0),
            fill_color=self._get_color("primary"),
        )
        self.add_textbox(
            title,
            Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7),
            font_size=28, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _add_section_title(self, title, subtitle=""):
        """章节标题页"""
        self.create_slide()
        self._add_bg()
        accent = self._get_color("accent")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Inches(10), Inches(2.0),
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.10)
        self.add_textbox(
            title,
            Inches(1), Inches(3.0), Inches(8), Inches(1.0),
            font_size=36, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self.add_textbox(
                subtitle,
                Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                font_size=18, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _add_thank_you(self):
        self.create_slide()
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [self._get_color("primary"), self._get_color("secondary")],
            angle=135,
        )
        self.add_textbox(
            "Thank You",
            Inches(1), Inches(2.2), Inches(8), Inches(1.2),
            font_size=48, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            "全新产品，即将上线",
            Inches(1), Inches(3.5), Inches(8), Inches(0.6),
            font_size=22, font_color=RGBColor(187, 222, 251),
            alignment=PP_ALIGN.CENTER,
        )
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.25),
            line_w, Pt(2),
            fill_color=RGBColor(187, 222, 251),
        )
        self.add_textbox(
            "敬请期待",
            Inches(1), Inches(4.2), Inches(8), Inches(0.5),
            font_size=16, font_color=RGBColor(100, 181, 246),
            alignment=PP_ALIGN.CENTER,
        )

    def _add_feature_card(self, left, top, width, height, icon_num, title, desc):
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
            fill_color=RGBColor(255, 255, 255),
        )
        # 图标
        icon_size = Inches(0.7)
        self.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.2), top + Inches(0.2),
            icon_size, icon_size,
            fill_color=self._get_color("accent"),
        )
        self.add_textbox(
            str(icon_num),
            left + Inches(0.2), top + Inches(0.28),
            icon_size, icon_size,
            font_size=18, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            title,
            left + Inches(1.0), top + Inches(0.2),
            width - Inches(1.2), Inches(0.35),
            font_size=14, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        self.add_textbox(
            desc,
            left + Inches(1.0), top + Inches(0.6),
            width - Inches(1.2), height - Inches(0.8),
            font_size=11, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.LEFT,
        )

    # ------------------------------------------------------------------ #
    #  各页面生成方法
    # ------------------------------------------------------------------ #

    def _gen_cover(self):
        """封面: 菱形展示"""
        self.create_slide()
        primary = self._get_color("primary")
        secondary = self._get_color("secondary")
        accent = self._get_color("accent")
        highlight = self._get_color("highlight")
        light = self._get_color("light")

        # 背景
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [primary, secondary], angle=120,
        )

        # 菱形装饰
        diamonds = [
            (Inches(0.5), Inches(0.3), Inches(2.5), Inches(2.5), 15, accent),
            (Inches(6.5), Inches(0.5), Inches(2.0), Inches(2.0), 0, highlight),
            (Inches(7.0), Inches(4.5), Inches(2.8), Inches(2.8), -10, light),
            (Inches(0.8), Inches(5.0), Inches(1.8), Inches(1.8), 30, accent),
            (Inches(4.0), Inches(5.5), Inches(1.5), Inches(1.5), 45, highlight),
        ]
        for left, top, w, h, rot, color in diamonds:
            shape = self.add_shape(
                MSO_SHAPE.DIAMOND, left, top, w, h,
                fill_color=color,
            )
            set_shape_transparency(shape, 0.55)
            if rot != 0:
                from lxml import etree
                xfrm = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    xfrm.set('rot', str(int(rot * 60000)))

        # 中心矩形遮罩
        mask_w = Inches(6.0)
        mask_h = Inches(2.8)
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            (Inches(10) - mask_w) // 2, (Inches(7.5) - mask_h) // 2,
            mask_w, mask_h,
            fill_color=RGBColor(0, 0, 0),
        )
        set_shape_transparency(self.slide.shapes[-1], 0.4)
        self.slide.shapes[-1].line.fill.background()

        # 标题
        self.add_textbox(
            self._title,
            Inches(2), Inches(2.5), Inches(6), Inches(1.2),
            font_size=44, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 分隔线
        line_w = Inches(2.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.75),
            line_w, Pt(2),
            fill_color=RGBColor(187, 222, 251),
        )
        # 副标题
        self.add_textbox(
            self._subtitle,
            Inches(2), Inches(4.0), Inches(6), Inches(0.8),
            font_size=20, font_color=RGBColor(187, 222, 251),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_directory(self):
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目 录")

        sections = [
            ("01", "痛点分析"),
            ("02", "解决方案"),
            ("03", "产品功能"),
            ("04", "产品演示"),
            ("05", "用户案例"),
            ("06", "价格方案"),
            ("07", "竞品对比"),
            ("08", "发布计划"),
        ]
        card_w = Inches(2.0)
        card_h = Inches(1.35)
        gap = Inches(0.3)
        cols = 4
        start_x = (Inches(10) - (cols * card_w + (cols - 1) * gap)) // 2
        start_y = Inches(1.5)

        for i, (num, title) in enumerate(sections):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            # 卡片
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(0.05),
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                num,
                x, y + Inches(0.2), card_w, Inches(0.4),
                font_size=22, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(0.7), card_w, Inches(0.35),
                font_size=13, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

    def _gen_pain_points(self):
        """痛点分析页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("痛点分析")

        pain_points = [
            ("效率低下", "传统工作流程繁琐\n重复性操作占用大量时间"),
            ("成本高昂", "人力成本持续攀升\nIT运维费用居高不下"),
            ("数据孤岛", "各部门数据不互通\n决策缺乏全局视角"),
            ("体验差", "界面复杂难用\n学习成本过高"),
        ]

        card_w = Inches(4.2)
        card_h = Inches(2.4)
        gap_x = Inches(0.4)
        gap_y = Inches(0.3)
        start_x = (Inches(10) - (2 * card_w + gap_x)) // 2
        start_y = Inches(1.5)

        for i, (title, desc) in enumerate(pain_points):
            row, col = divmod(i, 2)
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 左侧红色竖条
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, Inches(0.06), card_h,
                fill_color=RGBColor(244, 67, 54),
            )
            # 警告图标
            icon_size = Inches(0.6)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(0.3), y + Inches(0.3),
                icon_size, icon_size,
                fill_color=RGBColor(244, 67, 54),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.15)
            self.add_textbox(
                "!",
                x + Inches(0.3), y + Inches(0.35),
                icon_size, icon_size,
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标题
            self.add_textbox(
                title,
                x + Inches(1.0), y + Inches(0.3),
                card_w - Inches(1.2), Inches(0.35),
                font_size=16, font_color=RGBColor(244, 67, 54),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            # 描述
            self.add_textbox(
                desc,
                x + Inches(0.3), y + Inches(1.1),
                card_w - Inches(0.6), Inches(1.0),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

    def _gen_solution(self):
        """解决方案页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("解决方案")

        solutions = [
            ("AI智能引擎", "基于深度学习的\n智能分析与决策"),
            ("一体化平台", "打通数据孤岛\n统一管理入口"),
            ("极简设计", "零学习成本\n所见即所得"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.2)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (title, desc) in enumerate(solutions):
            x = start_x + i * (card_w + gap)
            y = Inches(1.4)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 渐变顶部
            self.add_gradient_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(1.2),
                [self._get_color("primary"), self._get_color("accent")],
                angle=90,
            )
            # 编号
            self.add_textbox(
                f"0{i+1}",
                x, y + Inches(0.15), card_w, Inches(0.5),
                font_size=28, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 标题
            self.add_textbox(
                title,
                x, y + Inches(1.5), card_w, Inches(0.5),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                desc,
                x + Inches(0.3), y + Inches(2.2), card_w - Inches(0.6), Inches(1.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _gen_product_feature(self, page_num=1):
        """产品功能页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"产品功能 - 第{page_num}部分")

        features = {
            1: [
                {"icon": "1", "title": "智能数据分析", "desc": "一键生成可视化报表\nAI自动发现数据洞察\n支持自定义仪表盘"},
                {"icon": "2", "title": "自动化工作流", "desc": "拖拽式流程编排\n支持500+自动化规则\n跨系统数据同步"},
            ],
            2: [
                {"icon": "3", "title": "团队协作中心", "desc": "实时多人协作编辑\n评论与审批流程\n任务分配与追踪"},
                {"icon": "4", "title": "安全与合规", "desc": "企业级数据加密\n权限精细控制\nGDPR合规保障"},
            ],
            3: [
                {"icon": "5", "title": "API开放平台", "desc": "RESTful API接口\nWebhook事件通知\n开发者友好文档"},
                {"icon": "6", "title": "移动端体验", "desc": "原生iOS/Android应用\n离线数据同步\n推送实时通知"},
            ],
            4: [
                {"icon": "7", "title": "智能推荐", "desc": "个性化内容推荐\n用户行为预测\nA/B测试支持"},
                {"icon": "8", "title": "多语言支持", "desc": "支持12种语言\n本地化适配\n文化敏感性处理"},
            ],
        }

        items = features.get(page_num, features[1])
        card_w = Inches(4.2)
        card_h = Inches(3.5)
        gap = Inches(0.5)
        start_x = (Inches(10) - (2 * card_w + gap)) // 2

        for i, item in enumerate(items):
            x = start_x + i * (card_w + gap)
            y = Inches(1.5)
            self._add_feature_card(x, y, card_w, card_h,
                                   item["icon"], item["title"], item["desc"])

    def _gen_demo(self):
        """产品演示页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("产品演示")

        # 模拟屏幕框
        screen_w = Inches(8.0)
        screen_h = Inches(4.5)
        screen_x = (Inches(10) - screen_w) // 2
        screen_y = Inches(1.5)

        # 屏幕边框
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            screen_x - Inches(0.1), screen_y - Inches(0.1),
            screen_w + Inches(0.2), screen_h + Inches(0.2),
            fill_color=RGBColor(33, 33, 33),
        )
        # 屏幕内容区
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            screen_x, screen_y, screen_w, screen_h,
            fill_color=RGBColor(240, 247, 255),
        )
        # 顶部工具栏
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            screen_x, screen_y, screen_w, Inches(0.4),
            fill_color=self._get_color("primary"),
        )
        # 圆点装饰
        for j, color in enumerate([RGBColor(255, 95, 87), RGBColor(255, 189, 46), RGBColor(39, 201, 63)]):
            self.add_shape(
                MSO_SHAPE.OVAL,
                screen_x + Inches(0.2 + j * 0.3), screen_y + Inches(0.08),
                Inches(0.2), Inches(0.2),
                fill_color=color,
            )

        # 模拟内容区
        sidebar_w = Inches(1.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            screen_x, screen_y + Inches(0.4), sidebar_w, screen_h - Inches(0.4),
            fill_color=self._get_color("primary"),
        )
        set_shape_transparency(self.slide.shapes[-1], 0.15)

        # 主内容区标题
        self.add_textbox(
            "产品演示界面",
            screen_x + sidebar_w + Inches(0.3), screen_y + Inches(0.6),
            Inches(5.0), Inches(0.4),
            font_size=16, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        # 模拟图表区
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            screen_x + sidebar_w + Inches(0.3), screen_y + Inches(1.2),
            Inches(3.5), Inches(2.0),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "数据可视化区域",
            screen_x + sidebar_w + Inches(0.3), screen_y + Inches(1.8),
            Inches(3.5), Inches(0.5),
            font_size=12, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.CENTER,
        )

        # 右侧面板
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            screen_x + sidebar_w + Inches(4.2), screen_y + Inches(0.6),
            Inches(2.0), Inches(3.3),
            fill_color=RGBColor(255, 255, 255),
        )

        # 底部说明
        self.add_textbox(
            "实时数据处理 | 智能分析 | 一键导出",
            Inches(1), Inches(6.3), Inches(8), Inches(0.4),
            font_size=13, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_user_case(self, page_num=1):
        """用户案例页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"用户案例 - 第{page_num}部分")

        cases = {
            1: [
                {"company": "某大型制造企业", "industry": "制造业",
                 "result": "效率提升 40%", "desc": "通过自动化工作流，将生产排程时间从3天缩短到4小时"},
                {"company": "某互联网公司", "industry": "互联网",
                 "result": "成本降低 30%", "desc": "统一数据平台，消除重复数据处理，节省人力成本"},
            ],
            2: [
                {"company": "某金融机构", "industry": "金融",
                 "result": "决策提速 50%", "desc": "AI实时数据分析，帮助管理层快速做出投资决策"},
                {"company": "某零售集团", "industry": "零售",
                 "result": "营收增长 25%", "desc": "精准用户画像+个性化推荐，提升转化率和复购率"},
            ],
        }

        items = cases.get(page_num, cases[1])
        card_w = Inches(8.5)
        card_h = Inches(2.2)
        start_y = Inches(1.5)

        for i, case in enumerate(items):
            y = start_y + i * (card_h + Inches(0.4))
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.75), y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 左侧色块
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.75), y, Inches(0.06), card_h,
                fill_color=self._get_color("accent"),
            )
            # 公司名
            self.add_textbox(
                case["company"],
                Inches(1.1), y + Inches(0.3), Inches(3.5), Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            # 行业标签
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.1), y + Inches(0.8), Inches(1.2), Inches(0.3),
                fill_color=self._get_color("light"),
            )
            self.add_textbox(
                case["industry"],
                Inches(1.1), y + Inches(0.8), Inches(1.2), Inches(0.3),
                font_size=9, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 描述
            self.add_textbox(
                case["desc"],
                Inches(1.1), y + Inches(1.3), Inches(4.5), Inches(0.6),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            # 成果高亮
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.5), y + Inches(0.4), Inches(2.5), Inches(1.4),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.1)
            self.add_textbox(
                case["result"],
                Inches(6.5), y + Inches(0.7), Inches(2.5), Inches(0.6),
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

    def _gen_pricing(self):
        """价格方案页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("价格方案")

        plans = [
            {"name": "基础版", "price": "免费", "period": "",
             "features": ["5个用户", "10GB存储", "基础分析", "邮件支持"],
             "highlight": False},
            {"name": "专业版", "price": "299", "period": "元/月",
             "features": ["50个用户", "100GB存储", "高级分析", "优先支持"],
             "highlight": True},
            {"name": "企业版", "price": "联系我们", "period": "",
             "features": ["不限用户", "无限存储", "定制功能", "专属顾问"],
             "highlight": False},
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, plan in enumerate(plans):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            # 卡片
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            if plan["highlight"]:
                # 高亮边框
                self.slide.shapes[-1].line.color.rgb = self._get_color("accent")
                self.slide.shapes[-1].line.width = Pt(2)
                # 推荐标签
                self.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x + Inches(0.8), y - Inches(0.2),
                    Inches(1.0), Inches(0.35),
                    fill_color=self._get_color("accent"),
                )
                self.add_textbox(
                    "推荐",
                    x + Inches(0.8), y - Inches(0.2),
                    Inches(1.0), Inches(0.35),
                    font_size=10, font_color=RGBColor(255, 255, 255),
                    bold=True, alignment=PP_ALIGN.CENTER,
                )
            else:
                self.slide.shapes[-1].line.fill.background()

            # 名称
            self.add_textbox(
                plan["name"],
                x, y + Inches(0.3), card_w, Inches(0.4),
                font_size=18, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 价格
            self.add_textbox(
                plan["price"],
                x, y + Inches(0.9), card_w, Inches(0.6),
                font_size=32, font_color=self._get_color("accent") if plan["highlight"] else self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            if plan["period"]:
                self.add_textbox(
                    plan["period"],
                    x, y + Inches(1.4), card_w, Inches(0.3),
                    font_size=11, font_color=self._get_color("text_secondary"),
                    alignment=PP_ALIGN.CENTER,
                )
            # 分隔线
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(0.3), y + Inches(1.8),
                card_w - Inches(0.6), Pt(1),
                fill_color=self._get_color("light"),
            )
            # 功能列表
            fy = y + Inches(2.0)
            for feat in plan["features"]:
                self.add_textbox(
                    f"  {feat}",
                    x + Inches(0.3), fy,
                    card_w - Inches(0.6), Inches(0.3),
                    font_size=11, font_color=self._get_color("text_secondary"),
                    alignment=PP_ALIGN.LEFT,
                )
                fy += Inches(0.35)

    def _gen_competitor(self):
        """竞品对比页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("竞品对比")

        headers = ["功能", "我们", "竞品A", "竞品B"]
        col_w = [Inches(2.0), Inches(2.3), Inches(2.3), Inches(2.3)]
        start_x = Inches(0.6)
        header_y = Inches(1.4)

        for i, (h, w) in enumerate(zip(headers, col_w)):
            x = start_x
            for j in range(i):
                x += col_w[j]
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, header_y, w, Inches(0.5),
                fill_color=self._get_color("primary") if i == 0 else self._get_color("accent"),
            )
            self.add_textbox(
                h, x, header_y, w, Inches(0.5),
                font_size=13, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

        rows = [
            ("AI分析", "内置AI引擎", "需额外付费", "无"),
            ("协作功能", "实时协作", "基础协作", "无"),
            ("API开放", "完全开放", "部分开放", "不支持"),
            ("数据安全", "企业级加密", "基础加密", "标准加密"),
            ("价格", "更具竞争力", "偏高", "较低但功能少"),
        ]

        for r, row in enumerate(rows):
            y = header_y + Inches(0.5 + r * 0.55)
            bg = RGBColor(240, 247, 255) if r % 2 == 0 else RGBColor(255, 255, 255)
            for i, (val, w) in enumerate(zip(row, col_w)):
                x = start_x
                for j in range(i):
                    x += col_w[j]
                self.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5), fill_color=bg)
                fc = self._get_color("text_primary") if i == 0 else self._get_color("text_secondary")
                self.add_textbox(
                    val, x, y, w, Inches(0.5),
                    font_size=11, font_color=fc,
                    bold=(i == 0), alignment=PP_ALIGN.CENTER,
                )

    def _gen_launch_plan(self):
        """发布计划页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("发布计划")

        phases = [
            ("Phase 1", "2025 Q3", "内测阶段", "核心功能开发\n种子用户邀请\nBug修复优化"),
            ("Phase 2", "2025 Q4", "公测阶段", "开放注册\n功能完善\n用户反馈收集"),
            ("Phase 3", "2026 Q1", "正式发布", "全渠道推广\n合作伙伴对接\n规模化运营"),
        ]

        # 时间轴
        line_y = Inches(3.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), line_y, Inches(8.4), Pt(3),
            fill_color=self._get_color("accent"),
        )

        node_w = Inches(2.5)
        total_w = len(phases) * node_w + (len(phases) - 1) * Inches(0.4)
        start_x = (Inches(10) - total_w) // 2

        for i, (phase, time, title, desc) in enumerate(phases):
            x = start_x + i * (node_w + Inches(0.4))
            # 节点
            circle_size = Inches(0.5)
            cx = x + node_w // 2 - circle_size // 2
            self.add_shape(
                MSO_SHAPE.OVAL,
                cx, line_y - circle_size // 2,
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                str(i + 1),
                cx, line_y - Inches(0.15),
                circle_size, circle_size,
                font_size=14, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

            # 卡片
            card_y = line_y - Inches(2.2) if i % 2 == 0 else line_y + Inches(0.6)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, card_y, node_w, Inches(1.8),
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_textbox(
                time,
                x, card_y + Inches(0.1), node_w, Inches(0.3),
                font_size=11, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, card_y + Inches(0.4), node_w, Inches(0.35),
                font_size=13, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                desc,
                x + Inches(0.15), card_y + Inches(0.8), node_w - Inches(0.3), Inches(0.8),
                font_size=9, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    # ------------------------------------------------------------------ #
    #  主生成方法
    # ------------------------------------------------------------------ #

    def generate(self):
        """生成完整20页产品发布会"""
        self._gen_cover()           # 1. 封面
        self._gen_directory()       # 2. 目录
        self._add_section_title("痛点分析", "Pain Points")  # 3. 章节标题
        self._gen_pain_points()     # 4. 痛点分析
        self._gen_solution()        # 5. 解决方案
        self._add_section_title("产品功能", "Product Features")  # 6. 章节标题
        self._gen_product_feature(1)  # 7. 产品功能1
        self._gen_product_feature(2)  # 8. 产品功能2
        self._gen_product_feature(3)  # 9. 产品功能3
        self._gen_product_feature(4)  # 10. 产品功能4
        self._add_section_title("产品演示与案例", "Demo & Cases")  # 11. 章节标题
        self._gen_demo()            # 12. 产品演示
        self._gen_user_case(1)      # 13. 用户案例1
        self._gen_user_case(2)      # 14. 用户案例2
        self._add_section_title("商务合作", "Business")  # 15. 章节标题
        self._gen_pricing()         # 16. 价格方案
        self._gen_competitor()      # 17. 竞品对比
        self._add_section_title("发布计划", "Launch Plan")  # 18. 章节标题
        self._gen_launch_plan()     # 19. 发布计划
        self._add_thank_you()       # 20. 致谢
        return self


if __name__ == "__main__":
    assembler = ProductLaunchAssembler()
    assembler.set_title("产品发布会").set_subtitle("全新产品 重磅发布")
    assembler.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_product_launch.pptx')
    assembler.save(output_path)
    print(f"产品发布会已生成: {output_path}")

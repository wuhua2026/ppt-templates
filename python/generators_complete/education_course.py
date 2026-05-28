"""教育课件完整演示文稿组装器

生成18页教育课件演示文稿：
封面 -> 目录 -> 课程介绍 -> 教学目标 -> 知识点(6) -> 案例分析(2) -> 课堂练习 -> 总结回顾 -> 致谢

默认使用绿色自然主题，极简渐变封面。
"""

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from python.base import TemplateGenerator
from python.utils import hex_to_rgb, set_shape_transparency


class EducationCourseAssembler(TemplateGenerator):
    """教育课件演示文稿组装器 (18 slides)"""

    DEFAULT_THEME = {
        "primary": "#1B5E20",
        "secondary": "#2E7D32",
        "accent": "#43A047",
        "highlight": "#81C784",
        "light": "#E8F5E9",
        "bg": "#FAFAFA",
        "text_primary": "#1B5E20",
        "text_secondary": "#455A64",
        "text_light": "#FFFFFF",
    }

    def __init__(self, theme=None):
        super().__init__(theme=theme)
        self._title = "教育课件"
        self._subtitle = "Education Course"
        self._data = {}
        self._theme = theme or self.DEFAULT_THEME

    def set_title(self, title):
        self._title = title
        return self

    def set_subtitle(self, subtitle):
        self._subtitle = subtitle
        return self

    def set_data(self, data):
        """设置演示文稿数据

        data: dict, 可包含以下键:
            - course_intro: 课程介绍
            - objectives: 教学目标
            - knowledge_points: 知识点列表
            - cases: 案例分析
            - exercises: 课堂练习
            - summary: 总结回顾
        """
        self._data = data
        return self

    def _get_color(self, key):
        return hex_to_rgb(getattr(self._theme, key, self._theme.primary))

    # ------------------------------------------------------------------ #
    #  通用辅助方法
    # ------------------------------------------------------------------ #

    def _add_bg(self, color_key="bg"):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            fill_color=self._get_color(color_key),
        )

    def _add_page_title_bar(self, title):
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(1.0),
            fill_color=self._get_color("primary"),
        )
        self.add_textbox(
            title,
            Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7),
            font_size=28, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

    def _add_section_title(self, title, subtitle=""):
        """章节标题页"""
        self.create_slide()
        self._add_bg()
        accent = self._get_color("accent")
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Inches(10), Inches(2.0),
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.10)
        self.add_textbox(
            title,
            Inches(1), Inches(3.0), Inches(8), Inches(1.0),
            font_size=36, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self.add_textbox(
                subtitle,
                Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                font_size=18, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

    def _add_thank_you(self):
        self.create_slide()
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [self._get_color("primary"), self._get_color("secondary")],
            angle=135,
        )
        self.add_textbox(
            "Thank You",
            Inches(1), Inches(2.2), Inches(8), Inches(1.2),
            font_size=48, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            "感谢聆听 祝学习愉快",
            Inches(1), Inches(3.5), Inches(8), Inches(0.8),
            font_size=22, font_color=RGBColor(129, 199, 132),
            alignment=PP_ALIGN.CENTER,
        )
        line_w = Inches(2)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.3),
            line_w, Pt(2),
            fill_color=RGBColor(129, 199, 132),
        )

    def _add_knowledge_card(self, left, top, width, height, num, title, points):
        """知识点卡片"""
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
            fill_color=RGBColor(255, 255, 255),
        )
        # 编号圆
        circle_size = Inches(0.5)
        self.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.2), top + Inches(0.2),
            circle_size, circle_size,
            fill_color=self._get_color("accent"),
        )
        self.add_textbox(
            str(num),
            left + Inches(0.2), top + Inches(0.25),
            circle_size, circle_size,
            font_size=16, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 标题
        self.add_textbox(
            title,
            left + Inches(0.8), top + Inches(0.2),
            width - Inches(1.0), Inches(0.35),
            font_size=14, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        # 要点
        y_pos = top + Inches(0.7)
        for point in points[:4]:
            self.add_textbox(
                f"  {point}",
                left + Inches(0.2), y_pos,
                width - Inches(0.4), Inches(0.25),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            y_pos += Inches(0.28)

    # ------------------------------------------------------------------ #
    #  各页面生成方法
    # ------------------------------------------------------------------ #

    def _gen_cover(self):
        """封面: 极简渐变"""
        self.create_slide()
        primary = self._get_color("primary")
        secondary = self._get_color("secondary")
        accent = self._get_color("accent")
        highlight = self._get_color("highlight")

        # 背景渐变
        self.add_gradient_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5),
            [primary, secondary], angle=45,
        )

        # 装饰圆
        circle_size = Inches(4.5)
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(10) - circle_size + Inches(0.5),
            Inches(-1.5),
            circle_size, circle_size,
            fill_color=accent,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.85)

        circle_size2 = Inches(3.0)
        self.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1.0),
            Inches(7.5) - circle_size2 + Inches(0.5),
            circle_size2, circle_size2,
            fill_color=highlight,
        )
        set_shape_transparency(self.slide.shapes[-1], 0.88)

        # 标题
        self.add_textbox(
            self._title,
            Inches(1.5), Inches(2.6), Inches(7.0), Inches(1.2),
            font_size=44, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # 分隔线
        line_w = Inches(2.5)
        self.add_shape(
            MSO_SHAPE.RECTANGLE,
            (Inches(10) - line_w) // 2, Inches(3.85),
            line_w, Pt(2),
            fill_color=RGBColor(129, 199, 132),
        )
        # 副标题
        self.add_textbox(
            self._subtitle,
            Inches(1.5), Inches(4.1), Inches(7.0), Inches(0.8),
            font_size=20, font_color=RGBColor(129, 199, 132),
            alignment=PP_ALIGN.CENTER,
        )

    def _gen_directory(self):
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("目 录")

        sections = [
            ("01", "课程介绍"),
            ("02", "教学目标"),
            ("03", "知识点讲解"),
            ("04", "案例分析"),
            ("05", "课堂练习"),
            ("06", "总结回顾"),
        ]
        card_w = Inches(2.6)
        card_h = Inches(1.5)
        gap = Inches(0.4)
        cols = 3
        start_x = (Inches(10) - (cols * card_w + (cols - 1) * gap)) // 2
        start_y = Inches(1.5)

        for i, (num, title) in enumerate(sections):
            row, col = divmod(i, cols)
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(0.06),
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                num,
                x, y + Inches(0.2), card_w, Inches(0.5),
                font_size=24, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(0.8), card_w, Inches(0.4),
                font_size=14, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )

    def _gen_course_intro(self):
        """课程介绍页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("课程介绍")

        # 左侧
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "课程概述",
            Inches(0.8), Inches(1.5), Inches(5.0), Inches(0.4),
            font_size=18, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        intro_lines = [
            "本课程旨在帮助学员系统掌握核心知识与技能。",
            "通过理论讲解与实践操作相结合的方式，",
            "让学员能够在实际工作中灵活应用所学内容。",
            "",
            "课程时长: 4课时 (共180分钟)",
            "难度级别: 中级",
            "适用人群: 有基础的进阶学习者",
        ]
        y = Inches(2.1)
        for line in intro_lines:
            self.add_textbox(
                line,
                Inches(0.8), y, Inches(5.0), Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            y += Inches(0.35)

        # 右侧 - 课程结构
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.3), Inches(1.3), Inches(3.2), Inches(5.5),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "课程结构",
            Inches(6.5), Inches(1.5), Inches(2.8), Inches(0.4),
            font_size=16, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        structure = [
            ("第1课时", "基础概念与原理"),
            ("第2课时", "核心知识点讲解"),
            ("第3课时", "案例分析与讨论"),
            ("第4课时", "实践练习与总结"),
        ]
        y = Inches(2.2)
        for time, desc in structure:
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.5), y, Inches(2.8), Inches(0.9),
                fill_color=self._get_color("light"),
            )
            self.add_textbox(
                time,
                Inches(6.6), y + Inches(0.1), Inches(2.6), Inches(0.3),
                font_size=11, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            self.add_textbox(
                desc,
                Inches(6.6), y + Inches(0.45), Inches(2.6), Inches(0.3),
                font_size=10, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            y += Inches(1.1)

    def _gen_objectives(self):
        """教学目标页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("教学目标")

        objectives = [
            ("知识目标", [
                "理解核心概念和基本原理",
                "掌握关键知识点的内在联系",
                "了解行业最新发展趋势",
            ]),
            ("能力目标", [
                "能够独立分析和解决问题",
                "熟练运用所学工具和方法",
                "具备实际项目操作能力",
            ]),
            ("素养目标", [
                "培养批判性思维能力",
                "提升团队协作与沟通能力",
                "建立终身学习的意识",
            ]),
        ]

        card_w = Inches(2.6)
        card_h = Inches(4.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (title, items) in enumerate(objectives):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 顶部色块
            self.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, card_w, Inches(1.0),
                fill_color=self._get_color("accent"),
            )
            set_shape_transparency(self.slide.shapes[-1], 0.1)
            self.add_textbox(
                f"0{i+1}",
                x, y + Inches(0.1), card_w, Inches(0.4),
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(0.5), card_w, Inches(0.35),
                font_size=13, font_color=RGBColor(255, 255, 255),
                alignment=PP_ALIGN.CENTER,
            )
            # 列表
            iy = y + Inches(1.3)
            for item in items:
                self.add_textbox(
                    f"  {item}",
                    x + Inches(0.2), iy, card_w - Inches(0.4), Inches(0.35),
                    font_size=11, font_color=self._get_color("text_secondary"),
                    alignment=PP_ALIGN.LEFT,
                )
                iy += Inches(0.45)

    def _gen_knowledge_point(self, page_num=1):
        """知识点页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"知识点 - 第{page_num}部分")

        knowledge_data = {
            1: ("核心概念", [
                {"title": "基本定义", "points": ["明确概念边界", "区分相似概念", "理解应用场景"]},
                {"title": "理论基础", "points": ["经典理论回顾", "现代发展演变", "关键假设条件"]},
            ]),
            2: ("关键原理", [
                {"title": "原理一", "points": ["核心逻辑链条", "数学推导过程", "直觉理解方法"]},
                {"title": "原理二", "points": ["实际应用场景", "边界条件分析", "常见误区提醒"]},
            ]),
            3: ("方法论", [
                {"title": "分析框架", "points": ["系统性思考方法", "分类分级策略", "优先级排序技巧"]},
                {"title": "工具使用", "points": ["常用工具介绍", "操作步骤演示", "最佳实践分享"]},
            ]),
            4: ("实操技巧", [
                {"title": "常见问题", "points": ["Top 10 高频问题", "解决方案模板", "避坑指南"]},
                {"title": "效率提升", "points": ["快捷操作技巧", "自动化方案", "协同工作流程"]},
            ]),
            5: ("进阶内容", [
                {"title": "深度分析", "points": ["高级分析方法", "数据驱动决策", "模型构建技巧"]},
                {"title": "行业应用", "points": ["跨行业案例", "定制化方案", "趋势预判方法"]},
            ]),
            6: ("前沿动态", [
                {"title": "技术趋势", "points": ["最新技术发展", "AI与自动化", "行业变革方向"]},
                {"title": "最佳实践", "points": ["标杆企业经验", "创新实践案例", "未来发展方向"]},
            ]),
        }

        section_title, items = knowledge_data.get(page_num, ("知识点", []))

        for i, item in enumerate(items):
            x = Inches(0.5) if i == 0 else Inches(5.0)
            self._add_knowledge_card(
                x, Inches(1.3), Inches(4.5), Inches(5.5),
                (page_num - 1) * 2 + i + 1, item["title"], item["points"]
            )

    def _gen_case_study(self, page_num=1):
        """案例分析页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar(f"案例分析 - 第{page_num}部分")

        cases = {
            1: {
                "title": "案例一: 某企业数字化转型",
                "background": "某传统制造企业面临效率低下、数据孤岛等问题",
                "approach": "采用分阶段数字化策略，先打通核心业务流程",
                "result": "生产效率提升35%，运营成本降低20%",
                "lessons": ["高层支持是关键", "先易后难的推进策略", "数据治理是基础"],
            },
            2: {
                "title": "案例二: 某互联网产品增长",
                "background": "新产品上线后用户增长缓慢，留存率低",
                "approach": "通过A/B测试优化用户路径，个性化推荐提升体验",
                "result": "月活用户增长200%，次日留存提升15个百分点",
                "lessons": ["数据驱动决策", "快速迭代验证", "用户反馈闭环"],
            },
        }

        case = cases.get(page_num, cases[1])

        # 案例标题
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.6),
            fill_color=self._get_color("accent"),
        )
        set_shape_transparency(self.slide.shapes[-1], 0.1)
        self.add_textbox(
            case["title"],
            Inches(0.8), Inches(1.35), Inches(8.4), Inches(0.5),
            font_size=16, font_color=RGBColor(255, 255, 255),
            bold=True, alignment=PP_ALIGN.LEFT,
        )

        sections = [
            ("背景", case["background"]),
            ("方法", case["approach"]),
            ("成果", case["result"]),
        ]

        y = Inches(2.1)
        for label, content in sections:
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), y, Inches(9.0), Inches(1.0),
                fill_color=RGBColor(255, 255, 255),
            )
            # 标签
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), y + Inches(0.2), Inches(0.8), Inches(0.35),
                fill_color=self._get_color("light"),
            )
            self.add_textbox(
                label,
                Inches(0.7), y + Inches(0.2), Inches(0.8), Inches(0.35),
                font_size=10, font_color=self._get_color("accent"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                content,
                Inches(1.7), y + Inches(0.25), Inches(7.6), Inches(0.5),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            y += Inches(1.15)

        # 启示
        self.add_textbox(
            "启示与思考",
            Inches(0.7), Inches(5.6), Inches(8.4), Inches(0.35),
            font_size=14, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        lesson_text = " | ".join(case["lessons"])
        self.add_textbox(
            lesson_text,
            Inches(0.7), Inches(6.0), Inches(8.4), Inches(0.5),
            font_size=11, font_color=self._get_color("text_secondary"),
            alignment=PP_ALIGN.LEFT,
        )

    def _gen_exercise(self):
        """课堂练习页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("课堂练习")

        exercises = [
            {"num": "1", "question": "请根据本节课内容，画出核心概念的思维导图。",
             "hint": "提示: 从中心概念出发，至少延伸3个分支"},
            {"num": "2", "question": "分析以下场景，运用所学方法提出解决方案。",
             "hint": "提示: 结合案例中的分析框架"},
            {"num": "3", "question": "小组讨论: 如何将本节课知识应用到实际工作中?",
             "hint": "提示: 每组15分钟讨论，5分钟汇报"},
        ]

        card_w = Inches(9.0)
        card_h = Inches(1.6)
        start_y = Inches(1.3)

        for i, ex in enumerate(exercises):
            y = start_y + i * (card_h + Inches(0.25))
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 编号
            circle_size = Inches(0.5)
            self.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.7), y + Inches(0.3),
                circle_size, circle_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                ex["num"],
                Inches(0.7), y + Inches(0.35),
                circle_size, circle_size,
                font_size=16, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            # 题目
            self.add_textbox(
                ex["question"],
                Inches(1.4), y + Inches(0.3), Inches(7.8), Inches(0.5),
                font_size=14, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.LEFT,
            )
            # 提示
            self.add_textbox(
                ex["hint"],
                Inches(1.4), y + Inches(0.9), Inches(7.8), Inches(0.4),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )

        # 注意事项
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.5), Inches(9.0), Inches(0.7),
            fill_color=self._get_color("light"),
        )
        self.add_textbox(
            "注意: 请在练习本上完成，完成后与同桌交流讨论，教师将进行点评。",
            Inches(0.8), Inches(6.6), Inches(8.4), Inches(0.4),
            font_size=11, font_color=self._get_color("text_primary"),
            alignment=PP_ALIGN.LEFT,
        )

    def _gen_summary(self):
        """总结回顾页"""
        self.create_slide()
        self._add_bg()
        self._add_page_title_bar("总结回顾")

        summary_items = [
            ("核心概念", "掌握了基本定义和理论基础\n理解了概念之间的内在联系"),
            ("方法工具", "学会了分析框架和方法论\n掌握了常用工具的使用技巧"),
            ("实践应用", "通过案例加深了理解\n为实际工作打下了基础"),
        ]

        card_w = Inches(2.6)
        card_h = Inches(3.5)
        gap = Inches(0.4)
        start_x = (Inches(10) - (3 * card_w + 2 * gap)) // 2

        for i, (title, desc) in enumerate(summary_items):
            x = start_x + i * (card_w + gap)
            y = Inches(1.3)
            self.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_w, card_h,
                fill_color=RGBColor(255, 255, 255),
            )
            # 图标圆
            icon_size = Inches(0.7)
            self.add_shape(
                MSO_SHAPE.OVAL,
                x + (card_w - icon_size) // 2, y + Inches(0.3),
                icon_size, icon_size,
                fill_color=self._get_color("accent"),
            )
            self.add_textbox(
                str(i + 1),
                x + (card_w - icon_size) // 2, y + Inches(0.38),
                icon_size, icon_size,
                font_size=20, font_color=RGBColor(255, 255, 255),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                title,
                x, y + Inches(1.2), card_w, Inches(0.4),
                font_size=16, font_color=self._get_color("text_primary"),
                bold=True, alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                desc,
                x + Inches(0.2), y + Inches(1.8), card_w - Inches(0.4), Inches(1.5),
                font_size=11, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.CENTER,
            )

        # 课后任务
        self.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(5.2), Inches(9.0), Inches(1.8),
            fill_color=RGBColor(255, 255, 255),
        )
        self.add_textbox(
            "课后任务",
            Inches(0.8), Inches(5.3), Inches(8.4), Inches(0.4),
            font_size=16, font_color=self._get_color("text_primary"),
            bold=True, alignment=PP_ALIGN.LEFT,
        )
        tasks = [
            "1. 复习本节课知识点，整理课堂笔记",
            "2. 完成课后练习题 (教材 P45-P48)",
            "3. 预习下节课内容: 高级应用与实战",
        ]
        ty = Inches(5.8)
        for task in tasks:
            self.add_textbox(
                task,
                Inches(0.8), ty, Inches(8.4), Inches(0.3),
                font_size=12, font_color=self._get_color("text_secondary"),
                alignment=PP_ALIGN.LEFT,
            )
            ty += Inches(0.35)

    # ------------------------------------------------------------------ #
    #  主生成方法
    # ------------------------------------------------------------------ #

    def generate(self):
        """生成完整18页教育课件"""
        self._gen_cover()              # 1. 封面
        self._gen_directory()          # 2. 目录
        self._add_section_title("课程介绍", "Course Introduction")  # 3. 章节标题
        self._gen_course_intro()       # 4. 课程介绍
        self._gen_objectives()         # 5. 教学目标
        self._add_section_title("知识点讲解", "Knowledge Points")  # 6. 章节标题
        self._gen_knowledge_point(1)   # 7. 知识点1
        self._gen_knowledge_point(2)   # 8. 知识点2
        self._gen_knowledge_point(3)   # 9. 知识点3
        self._gen_knowledge_point(4)   # 10. 知识点4
        self._gen_knowledge_point(5)   # 11. 知识点5
        self._gen_knowledge_point(6)   # 12. 知识点6
        self._add_section_title("案例与练习", "Cases & Practice")  # 13. 章节标题
        self._gen_case_study(1)        # 14. 案例分析1
        self._gen_case_study(2)        # 15. 案例分析2
        self._gen_exercise()           # 16. 课堂练习
        self._gen_summary()            # 17. 总结回顾
        self._add_thank_you()          # 18. 致谢
        return self


if __name__ == "__main__":
    assembler = EducationCourseAssembler()
    assembler.set_title("教育课件").set_subtitle("数字化转型实战课程")
    assembler.generate()
    output_path = os.path.join(os.path.dirname(__file__), '..', '..', 'output_education_course.pptx')
    assembler.save(output_path)
    print(f"教育课件已生成: {output_path}")
